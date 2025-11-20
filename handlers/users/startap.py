import os
import asyncio
import sqlite3
import logging
import json
from datetime import datetime
from typing import Optional, Dict, List
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, InputFile, ContentType
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters.state import State, StatesGroup
from loader import bot, dp
from aiogram import types

from environs import Env

# environs kutubxonasidan foydalanish
env = Env()
env.read_env()

# .env fayl ichidan quyidagilarni o'qiymiz
OPENAI_API_KEY = env.str("OPENAI_API_KEY")  # Bot token

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# ==================== KONFIGURATSIYA ====================
USE_OPENAI = True  # Yoqilgan
ADMIN_ID = 736290914
# ADMIN_ID = 1879114908

CURRENCY = "so'm"
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
ALLOWED_FILE_TYPES = ['image/jpeg', 'image/png', 'application/pdf']

# Karta ma'lumotlari
CARD_NUMBER = "4073420066945407"
CARD_HOLDER = "Boburjon Astanov"

# Video file_id
VIDEO_FILE_ID = "BAACAgIAAxkBAAIEbWkdva_dB9kNCWcb8DmmDsdGo6AnAALoEgACpRlpSY0jnmDrEDCoNgQ"

# Logging sozlash
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# OpenAI client
openai_client = None
if USE_OPENAI and OpenAI and OPENAI_API_KEY:
    openai_client = OpenAI(api_key=OPENAI_API_KEY)

# ==================== SAVOLLAR RO'YXATI ====================
QUESTIONS = [
    "1️⃣ Ismingiz?",
    "2️⃣ Loyiha nomi?",
    "3️⃣ Loyiha tavsifi (qisqacha, 2-3 jumla)?",
    "4️⃣ Qanday muammoni hal qilasiz?",
    "5️⃣ Sizning yechimingiz?",
    "6️⃣ Maqsadli auditoriya kimlar?",
    "7️⃣ Biznes model (qanday daromad olasiz)?",
    "8️⃣ Asosiy raqobatchilaringiz?",
    "9️⃣ Sizning ustunligingiz (raqobatchilardan farqi)?",
    "🔟 Moliyaviy prognoz (keyingi 1 yil)?",
]


# ==================== DATABASE ====================
class Database:
    def __init__(self, db_path='pitch_bot.db'):
        self.db_path = db_path
        self.init_db()

    def init_db(self):
        """Ma'lumotlar bazasini yaratish"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                user_id INTEGER PRIMARY KEY,
                username TEXT,
                full_name TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS orders (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                package TEXT,
                status TEXT,
                answers TEXT,
                receipt_file_id TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users(user_id)
            )
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS admin_logs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                admin_id INTEGER,
                action TEXT,
                user_id INTEGER,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        conn.commit()
        conn.close()

    def save_user(self, user_id: int, username: str, full_name: str):
        """Foydalanuvchini saqlash"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT OR REPLACE INTO users (user_id, username, full_name)
            VALUES (?, ?, ?)
        ''', (user_id, username, full_name))
        conn.commit()
        conn.close()

    def create_order(self, user_id: int, answers: List[str], package: str) -> int:
        """Yangi buyurtma yaratish"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO orders (user_id, status, answers, package)
            VALUES (?, ?, ?, ?)
        ''', (user_id, 'awaiting_payment', '|||'.join(answers), package))
        order_id = cursor.lastrowid
        conn.commit()
        conn.close()
        return order_id

    def update_order(self, user_id: int, **kwargs):
        """Buyurtmani yangilash"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        set_clause = ', '.join([f"{k} = ?" for k in kwargs.keys()])
        values = list(kwargs.values()) + [user_id]
        cursor.execute(f'''
            UPDATE orders 
            SET {set_clause}, updated_at = CURRENT_TIMESTAMP
            WHERE user_id = ? AND status != 'completed'
        ''', values)
        conn.commit()
        conn.close()

    def get_order(self, user_id: int) -> Optional[Dict]:
        """Oxirgi buyurtmani olish"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute('''
            SELECT id, user_id, package, status, answers, receipt_file_id, created_at
            FROM orders
            WHERE user_id = ?
            ORDER BY created_at DESC
            LIMIT 1
        ''', (user_id,))
        row = cursor.fetchone()
        conn.close()
        if row:
            return {
                'id': row[0],
                'user_id': row[1],
                'package': row[2],
                'status': row[3],
                'answers': row[4].split('|||') if row[4] else [],
                'receipt_file_id': row[5],
                'created_at': row[6]
            }
        return None

    def log_admin_action(self, admin_id: int, action: str, user_id: int):
        """Admin harakatini loglash"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO admin_logs (admin_id, action, user_id)
            VALUES (?, ?, ?)
        ''', (admin_id, action, user_id))
        conn.commit()
        conn.close()


db = Database()


# ==================== STATES ====================
class PitchStates(StatesGroup):
    waiting_for_answer = State()
    choosing_package = State()
    waiting_for_receipt = State()


# ==================== KLAVIATURALAR ====================
def start_keyboard():
    kb = InlineKeyboardMarkup(row_width=2)
    kb.add(
        InlineKeyboardButton("✅ Boshlash", callback_data="start_yes"),
        InlineKeyboardButton("❌ Keyinroq", callback_data="start_no")
    )
    kb.add(
        InlineKeyboardButton("📹 Video ko'rish", callback_data="watch_video")
    )
    return kb


def package_keyboard():
    kb = InlineKeyboardMarkup(row_width=2)
    kb.add(
        InlineKeyboardButton(f"🔹 Oddiy — 10,000 {CURRENCY}", callback_data="package_simple"),
        InlineKeyboardButton(f"🔸 Pro — 20,000 {CURRENCY}", callback_data="package_pro")
    )
    kb.add(InlineKeyboardButton("❌ Bekor qilish", callback_data="cancel"))
    return kb


def admin_action_kb(user_id: int):
    kb = InlineKeyboardMarkup(row_width=2)
    kb.add(
        InlineKeyboardButton("✅ Tasdiqlash", callback_data=f"admin_approve:{user_id}"),
        InlineKeyboardButton("❌ Rad etish", callback_data=f"admin_reject:{user_id}")
    )
    kb.add(InlineKeyboardButton("👤 Ma'lumot", callback_data=f"admin_info:{user_id}"))
    return kb


def cancel_keyboard():
    kb = InlineKeyboardMarkup(row_width=1)
    kb.add(InlineKeyboardButton("❌ Bekor qilish", callback_data="cancel"))
    return kb


# ==================== HELPER FUNKSIYALAR ====================
async def validate_file(message: types.Message) -> tuple:
    """File validatsiyasi"""
    if message.content_type == ContentType.PHOTO:
        file_id = message.photo[-1].file_id
        file_info = await bot.get_file(file_id)
        file_size = file_info.file_size
        mime_type = 'image/jpeg'
    elif message.content_type == ContentType.DOCUMENT:
        file_id = message.document.file_id
        file_size = message.document.file_size
        mime_type = message.document.mime_type
    else:
        return False, "Faqat rasm yoki PDF yuklang", None

    if file_size > MAX_FILE_SIZE:
        max_size_mb = MAX_FILE_SIZE // 1024 // 1024
        return False, f"Fayl hajmi {max_size_mb} MB dan oshmasin", None

    if mime_type not in ALLOWED_FILE_TYPES:
        return False, "Faqat JPG, PNG yoki PDF formatda yuboring", None

    return True, "OK", file_id


def format_order_info(order: Dict) -> str:
    """Buyurtma ma'lumotlarini formatlash"""
    status_emoji = {
        'awaiting_payment': '💳',
        'receipt_sent': '⏳',
        'approved': '✅',
        'rejected': '❌',
        'completed': '🎉'
    }
    status_text = {
        'awaiting_payment': "To'lov kutilmoqda",
        'receipt_sent': 'Chek yuborildi',
        'approved': 'Tasdiqlandi',
        'rejected': 'Rad etildi',
        'completed': 'Yakunlandi'
    }
    package_text = 'Oddiy' if order.get('package') == 'simple' else 'Professional'

    status = order['status']
    emoji = status_emoji.get(status, '❓')
    text = status_text.get(status, 'Noma\'lum')
    package = package_text if order.get('package') else 'Tanlanmagan'
    created = order['created_at'][:16]

    result = f"{emoji} Holat: {text}\n"
    result += f"📦 Paket: {package}\n"
    result += f"📅 Yaratilgan: {created}"

    return result


# ==================== YANGILANGAN AI FUNKSIYALAR - BATAFSIL JAVOBLAR ====================
async def generate_market_analysis(project_info: str, target_audience: str, package: str) -> Dict:
    """Bozor tahlillari va statistik ma'lumotlar yaratish"""

    if not USE_OPENAI or not openai_client:
        return {
            'tam': "100 mln dollar",
            'sam': "30 mln dollar",
            'som': "5 mln dollar",
            'growth_rate': "25% yillik",
            'trends': "• Raqamlashtirish tendensiyasi\n• Mobil yechimlar talabi",
            'segments': "• B2B: 60%\n• B2C: 40%"
        }

    model = "gpt-4" if package == "pro" else "gpt-3.5-turbo"

    prompt = f"""
    Quyidagi loyiha uchun bozor tahlillarini O'ZBEK TILIDA yarating:
    Loyiha: {project_info}
    Maqsadli auditoriya: {target_audience}

    Quyidagilarni batafsil kiriting:
    1. TAM (Total Addressable Market) - umumiy bozor hajmi
    2. SAM (Serviceable Addressable Market) - xizmat ko'rsatish mumkin bo'lgan bozor
    3. SOM (Serviceable Obtainable Market) - erishish mumkin bo'lgan bozor ulushi
    4. Yillik o'sish sur'ati
    5. Bozor trendlari
    6. Mijozlar segmentatsiyasi

    JSON formatida qaytaring.
    """

    try:
        response = await asyncio.to_thread(
            lambda: openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "Siz bozor tahlili mutaxassisisiz. O'zbek tilida javob bering."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=1000,
                temperature=0.7,
                response_format={"type": "json_object"}
            )
        )

        return json.loads(response.choices[0].message.content)
    except Exception as e:
        logger.error(f"Market analysis failed: {e}")
        return {
            'tam': "100 mln dollar",
            'sam': "30 mln dollar",
            'som': "5 mln dollar",
            'growth_rate': "25% yillik",
            'trends': "• Raqamlashtirish tendensiyasi",
            'segments': "• B2B va B2C segmentlari"
        }


async def create_professional_pitch_content(answers: List[str], package: str) -> Dict:
    """AI bilan BATAFSIL va TO'LIQ pitch content yaratish"""

    if not USE_OPENAI or not openai_client:
        # Fallback content
        return generate_fallback_content(answers)

    model = "gpt-4" if package == "pro" else "gpt-3.5-turbo"

    # Avval bozor tahlillarini olish
    market_data = await generate_market_analysis(
        answers[1] if len(answers) > 1 else "",
        answers[5] if len(answers) > 5 else "",
        package
    )

    prompt = f"""
    Siz O'zbekistondagi eng yaxshi pitch deck mutaxassisisiz. 
    Quyidagi ma'lumotlar asosida JUDA BATAFSIL, TO'LIQ va PROFESSIONAL pitch content yarating.

    HAR BIR BO'LIM KAMIDA 5-7 TA TO'LIQ JUMLA BO'LISHI KERAK!

    STARTUP MA'LUMOTLARI:
    Asoschi: {answers[0] if len(answers) > 0 else ""}
    Loyiha: {answers[1] if len(answers) > 1 else ""}
    Tavsif: {answers[2] if len(answers) > 2 else ""}
    Muammo: {answers[3] if len(answers) > 3 else ""}
    Yechim: {answers[4] if len(answers) > 4 else ""}
    Maqsadli auditoriya: {answers[5] if len(answers) > 5 else ""}
    Biznes model: {answers[6] if len(answers) > 6 else ""}
    Raqobatchilar: {answers[7] if len(answers) > 7 else ""}
    Ustunliklar: {answers[8] if len(answers) > 8 else ""}
    Moliyaviy prognoz: {answers[9] if len(answers) > 9 else ""}

    BOZOR MA'LUMOTLARI:
    {json.dumps(market_data, ensure_ascii=False, indent=2)}

    O'ZBEK TILIDA JSON formatida JUDA BATAFSIL qaytaring:
    {{
      "project_name": "professional loyiha nomi",
      "author": "to'liq ism",
      "tagline": "ta'sirchan shior (8-10 so'z)",

      "problem_title": "BOZORDAGI MUAMMO",
      "problem": "• [BATAFSIL] Birinchi muammo: [kamida 2-3 jumla tushuntirish]\\n• [BATAFSIL] Ikkinchi muammo: [kamida 2-3 jumla tushuntirish]\\n• [BATAFSIL] Uchinchi muammo: [kamida 2-3 jumla tushuntirish]\\n• [BATAFSIL] To'rtinchi muammo: [kamida 2-3 jumla tushuntirish]\\n• [STATISTIKA] Muammo ta'sir ko'rsatayotgan odamlar soni va iqtisodiy zarar",

      "solution_title": "BIZNING INNOVATSION YECHIMIMIZ",  
      "solution": "• [ASOSIY YECHIM] Bizning platformamiz/mahsulotimiz qanday ishlaydi: [3-4 jumla batafsil]\\n• [TEXNOLOGIYA] Qanday texnologiyalar ishlatiladi: [2-3 jumla]\\n• [JARAYON] Mijoz uchun qanday qadamlar: [3-4 jumla]\\n• [NATIJA] Kutilayotgan natijalar va foyda: [2-3 jumla]\\n• [VAQT] Amalga oshirish vaqti va bosqichlari: [2-3 jumla]",

      "market_title": "BOZOR IMKONIYATLARI VA TAHLIL",
      "market": "📊 BOZOR HAJMI:\\n• Umumiy bozor (TAM): {market_data.get('tam', '')} - [batafsil tushuntirish]\\n• Mavjud bozor (SAM): {market_data.get('sam', '')} - [batafsil tushuntirish]\\n• Erishish mumkin (SOM): {market_data.get('som', '')} - [batafsil tushuntirish]\\n\\n📈 O'SISH SUR'ATI:\\n• Yillik o'sish: {market_data.get('growth_rate', '')}\\n• 3 yillik prognoz: [batafsil]\\n\\n🎯 MAQSADLI SEGMENTLAR:\\n[har bir segment uchun 2-3 jumla batafsil tavsif]",

      "business_title": "DAROMAD MODELI VA MONETIZATSIYA",
      "business_model": "💰 ASOSIY DAROMAD MANBALARI:\\n• [MODEL 1]: [batafsil tushuntirish, narxlar, misollar - 3-4 jumla]\\n• [MODEL 2]: [batafsil tushuntirish, narxlar, misollar - 3-4 jumla]\\n• [MODEL 3]: [batafsil tushuntirish, narxlar, misollar - 3-4 jumla]\\n\\n💳 NARXLASH STRATEGIYASI:\\n• [narx segmentlari va paketlar batafsil]\\n\\n📊 DAROMAD PROGNOZI:\\n• 1-oy: [summa va tushuntirish]\\n• 6-oy: [summa va tushuntirish]\\n• 1-yil: [summa va tushuntirish]",

      "competition_title": "RAQOBAT MUHITI VA TAHLILI", 
      "competition": "🏆 ASOSIY RAQOBATCHILAR:\\n• [RAQOBATCHI 1]: [kim, nima qiladi, kuchli/zaif tomonlari - 3-4 jumla]\\n• [RAQOBATCHI 2]: [kim, nima qiladi, kuchli/zaif tomonlari - 3-4 jumla]\\n• [RAQOBATCHI 3]: [kim, nima qiladi, kuchli/zaif tomonlari - 3-4 jumla]\\n\\n📊 BOZOR ULUSHLARI:\\n[har bir raqobatchining ulushi va tahlil]\\n\\n🎯 BIZNING POZITSIYAMIZ:\\n[bozordagi o'rni va strategiya - 3-4 jumla]",

      "advantage_title": "BIZNING YAGONA USTUNLIKLARIMIZ",
      "advantage": "⭐ [USTUNLIK 1 NOMI]:\\n[Bu ustunlik nimadan iborat, qanday amalga oshiriladi, mijozga qanday foyda - 4-5 jumla batafsil]\\n\\n⭐ [USTUNLIK 2 NOMI]:\\n[Bu ustunlik nimadan iborat, qanday amalga oshiriladi, mijozga qanday foyda - 4-5 jumla batafsil]\\n\\n⭐ [USTUNLIK 3 NOMI]:\\n[Bu ustunlik nimadan iborat, qanday amalga oshiriladi, mijozga qanday foyda - 4-5 jumla batafsil]\\n\\n⭐ [USTUNLIK 4 NOMI]:\\n[Bu ustunlik nimadan iborat, qanday amalga oshiriladi, mijozga qanday foyda - 4-5 jumla batafsil]",

      "financials_title": "MOLIYAVIY PROGNOZLAR VA KO'RSATKICHLAR",
      "financials": "📊 ASOSIY MOLIYAVIY KO'RSATKICHLAR:\\n\\n💰 DAROMAD PROGNOZI:\\n• 1-chorak: [summa] so'm - [tushuntirish]\\n• 2-chorak: [summa] so'm - [tushuntirish]\\n• 3-chorak: [summa] so'm - [tushuntirish]\\n• 4-chorak: [summa] so'm - [tushuntirish]\\n• YIL YAKUNIDA: [umumiy summa] so'm\\n\\n📈 ASOSIY METRIKALAR:\\n• CAC (mijoz jalb qilish narxi): [summa va tushuntirish]\\n• LTV (mijoz hayotiy qiymati): [summa va tushuntirish]\\n• Gross Margin: [foiz va tushuntirish]\\n• Break-even nuqtasi: [vaqt va tushuntirish]\\n\\n💸 INVESTITSIYA ZARURIYATI:\\n• Zarur summa: [summa] so'm\\n• Sarflanish yo'nalishlari: [batafsil]\\n• ROI prognozi: [foiz va vaqt]",

      "team_title": "JAMOA VA TAJRIBA",
      "team": "👥 ASOSIY JAMOA A'ZOLARI:\\n• [ISM 1] - [lavozim]: [tajriba va ko'nikmalar - 2-3 jumla]\\n• [ISM 2] - [lavozim]: [tajriba va ko'nikmalar - 2-3 jumla]\\n• [ISM 3] - [lavozim]: [tajriba va ko'nikmalar - 2-3 jumla]\\n\\n🏆 JAMOA YUTUQLARI:\\n[oldingi loyihalar, natijalar - 3-4 jumla]",

      "milestones_title": "YO'L XARITASI VA MUHIM BOSQICHLAR",
      "milestones": "🗓️ RIVOJLANISH BOSQICHLARI:\\n\\n✅ 0-3 OY:\\n• [Vazifa 1]: [batafsil tushuntirish]\\n• [Vazifa 2]: [batafsil tushuntirish]\\n• [Vazifa 3]: [batafsil tushuntirish]\\n\\n✅ 3-6 OY:\\n• [Vazifa 1]: [batafsil tushuntirish]\\n• [Vazifa 2]: [batafsil tushuntirish]\\n• [Vazifa 3]: [batafsil tushuntirish]\\n\\n✅ 6-12 OY:\\n• [Vazifa 1]: [batafsil tushuntirish]\\n• [Vazifa 2]: [batafsil tushuntirish]\\n• [Vazifa 3]: [batafsil tushuntirish]",

      "cta": "Keling, birgalikda O'zbekiston bozorida yangi standartlar o'rnatamiz! Bizning loyihamizga qo'shiling va kelajakni birgalikda quramiz! 🚀"
    }}

    ESLATMA: HAR BIR BO'LIM JUDA BATAFSIL VA TO'LIQ BO'LISHI KERAK!
    MINIMUM 5-7 TA TO'LIQ JUMLA HAR BIR ASOSIY PUNKTDA!
    """

    try:
        response = await asyncio.to_thread(
            lambda: openai_client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": "Siz O'zbekistondagi eng tajribali pitch deck mutaxassisisiz. Juda batafsil, to'liq va professional content yarating. Har bir bo'lim ko'p ma'lumot bilan to'ldirilishi kerak."
                    },
                    {"role": "user", "content": prompt}
                ],
                max_tokens=4000,  # Maksimal token
                temperature=0.8,
                response_format={"type": "json_object"}
            )
        )

        content = json.loads(response.choices[0].message.content)
        return content

    except Exception as e:
        logger.error(f"AI content creation failed: {e}")
        return generate_fallback_content(answers)


def generate_fallback_content(answers: List[str]) -> Dict:
    """Fallback content - AI ishlamasa"""
    return {
        'project_name': answers[1] if len(answers) > 1 else "Innovatsion Loyiha",
        'author': answers[0] if len(answers) > 0 else "Tadbirkor",
        'tagline': "Kelajakni birgalikda quramiz - innovatsiya orqali",
        'problem_title': "BOZORDAGI MUAMMO",
        'problem': f"""• Asosiy muammo: {answers[3] if len(answers) > 3 else 'Bozordagi samarasizlik'}
        Hozirgi kunda ko'plab kompaniyalar va shaxslar ushbu muammo bilan kurashmoqda.
        • Ikkinchi muammo: An'anaviy yechimlar samarasiz va qimmat
        Bu mijozlar uchun katta to'siq yaratmoqda va rivojlanishni sekinlashtirmoqda.
        • Uchinchi muammo: Texnologik yechimlar mavjud emas
        Bozorda hali bu muammoni to'liq hal qiladigan yechim yo'q.""",
        'solution_title': "BIZNING YECHIMIMIZ",
        'solution': f"""• Bizning yechim: {answers[4] if len(answers) > 4 else 'Innovatsion platforma'}
        Zamonaviy texnologiyalardan foydalangan holda muammoni hal qilamiz.
        • Qanday ishlaydi: Oddiy va qulay interfeys orqali
        Foydalanuvchilar bir necha daqiqada natijaga erisha oladilar.
        • Afzalliklari: Tez, arzon va samarali
        An'anaviy yechimlardan 10 barobar tezroq va 5 barobar arzonroq.""",
        'market_title': "BOZOR TAHLILI",
        'market': f"""📊 BOZOR HAJMI:
        • TAM: 500 mln dollar - O'zbekiston va Markaziy Osiyo bozori
        • SAM: 150 mln dollar - Bizning xizmat ko'rsata oladigan segment
        • SOM: 30 mln dollar - Birinchi 3 yilda erisha oladigan ulush

        🎯 MAQSADLI AUDITORIYA: {answers[5] if len(answers) > 5 else 'B2B va B2C segmentlari'}
        Asosiy mijozlarimiz 25-45 yosh oralig'idagi tadbirkorlar va kompaniyalar.""",
        'business_title': "BIZNES MODEL",
        'business_model': f"""💰 DAROMAD MODELI: {answers[6] if len(answers) > 6 else 'SaaS subscription'}
        • Oylik obuna: 50,000 - 500,000 so'm
        • Yillik paketlar: 20% chegirma bilan
        • Korporativ yechimlar: Individual narxlash

        📈 PROGNOZ:
        • 1-yil: 500 ta mijoz, 300 mln so'm daromad
        • 2-yil: 2000 ta mijoz, 1.5 mlrd so'm daromad
        • 3-yil: 10,000 ta mijoz, 8 mlrd so'm daromad""",
        'competition_title': "RAQOBAT",
        'competition': f"""🏆 RAQOBATCHILAR: {answers[7] if len(answers) > 7 else 'Mahalliy va xalqaro kompaniyalar'}
        • Asosiy raqobatchi: An'anaviy yechimlar
        Ular qimmat va sekin, texnologik jihatdan eskirgan
        • Ikkinchi raqobatchi: Xalqaro platformalar
        Mahalliy bozorga moslashmagan, qimmat va murakkab

        Bizning ustunligimiz: Mahalliy bozorni chuqur tushunish va zamonaviy texnologiya""",
        'advantage_title': "USTUNLIKLAR",
        'advantage': f"""⭐ BIZNING USTUNLIKLARIMIZ: {answers[8] if len(answers) > 8 else 'Yagona mahalliy yechim'}

        1. TEXNOLOGIK USTUNLIK:
        Sun'iy intellekt va avtomatlashtirish orqali 10x tezroq natijalar

        2. NARX USTUNLIGI:
        Raqobatchilarga nisbatan 50% arzonroq, lekin sifatli xizmat

        3. MAHALLIY EKSPERTIZA:
        O'zbek bozorini mukammal tushunish va mahalliy til/madaniyatga moslashgan

        4. MIJOZLARGA YONDASHUV:
        24/7 qo'llab-quvvatlash va shaxsiy yondashuv har bir mijozga""",
        'financials_title': "MOLIYAVIY KO'RSATKICHLAR",
        'financials': f"""📊 MOLIYAVIY PROGNOZ: {answers[9] if len(answers) > 9 else 'Ijobiy'}

        💰 DAROMAD:
        • 1-chorak: 50 mln so'm
        • 2-chorak: 150 mln so'm
        • 3-chorak: 300 mln so'm
        • 4-chorak: 500 mln so'm

        📈 ASOSIY METRIKALAR:
        • Gross Margin: 75%
        • CAC: 100,000 so'm
        • LTV: 2,000,000 so'm
        • Break-even: 18 oy

        💸 INVESTITSIYA:
        • Zarur: 500,000 USD
        • ROI: 300% (3 yilda)""",
        'team_title': "JAMOA",
        'team': """👥 PROFESSIONAL JAMOA:
        • CEO: 10+ yil tajriba texnologiya sohasida
        • CTO: 8+ yil dasturlash va AI tajribasi
        • CMO: 7+ yil marketing va sotish tajribasi
        • CFO: 12+ yil moliya va investitsiya tajribasi""",
        'milestones_title': "YO'L XARITASI",
        'milestones': """🗓️ MUHIM BOSQICHLAR:

        Q1 2024:
        • MVP ishga tushirish ✅
        • Birinchi 100 mijoz ✅
        • Seed investitsiya jalb qilish

        Q2-Q3 2024:
        • Mahsulotni kengaytirish
        • 500 mijozga yetish
        • Series A tayyorlash

        Q4 2024 - 2025:
        • Regional kengayish
        • 2000+ mijoz
        • Break-even ga chiqish""",
        'cta': "Keling, birgalikda O'zbekiston bozorida innovatsiya yaratamiz! 🚀"
    }


# ==================== MUKAMMAL PPTX YARATISH ====================
async def create_stunning_pitch_deck(user_id: int, answers: List[str], package: str) -> str:
    """Professional PPTX yaratish - maksimal content bilan"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR

    # AI orqali batafsil content olish
    content = await create_professional_pitch_content(answers, package)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Professional ranglar
    COLORS = {
        'primary': RGBColor(25, 42, 86),  # To'q ko'k
        'secondary': RGBColor(41, 128, 185),  # Och ko'k
        'accent': RGBColor(39, 174, 96),  # Yashil
        'danger': RGBColor(192, 57, 43),  # Qizil
        'warning': RGBColor(243, 156, 18),  # Sariq
        'purple': RGBColor(142, 68, 173),  # Binafsha
        'dark': RGBColor(44, 62, 80),  # To'q
        'light': RGBColor(236, 240, 241),  # Och kulrang
        'white': RGBColor(255, 255, 255),  # Oq
        'gray': RGBColor(149, 165, 166)  # Kulrang
    }

    def add_gradient_background(slide, color1, color2):
        """Gradient fon qo'shish"""
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = 45
        stops = fill.gradient_stops
        stops[0].color.rgb = color1
        stops[0].position = 0.0
        stops[1].color.rgb = color2
        stops[1].position = 1.0

    def add_content_slide(title, content_text, header_color):
        """Universal content slide yaratish"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['white']

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(32)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.text = content_text
        tf.word_wrap = True

        for p in tf.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(14)  # Kichikroq shrift ko'proq matn uchun
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0

        return slide

    # ==================== 1. BOSH SAHIFA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['primary'], COLORS['purple'])

    # Loyiha nomi
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.text = content.get('project_name', 'LOYIHA').upper()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(54)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True

    # Shior
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    tf = tagline_box.text_frame
    tf.text = content.get('tagline', '')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri Light"
    p.font.size = Pt(26)
    p.font.color.rgb = COLORS['light']
    p.font.italic = True

    # Taqdimotchi
    author_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
    tf = author_box.text_frame
    tf.text = f"Taqdim etmoqda: {content.get('author', '')}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['light']

    # ==================== 2. MUAMMO ====================
    if content.get('problem'):
        add_content_slide(
            f"🔥 {content.get('problem_title', 'MUAMMO')}",
            content.get('problem', ''),
            COLORS['danger']
        )

    # ==================== 3. YECHIM ====================
    if content.get('solution'):
        add_content_slide(
            f"💡 {content.get('solution_title', 'YECHIM')}",
            content.get('solution', ''),
            COLORS['accent']
        )

    # ==================== 4. BOZOR ====================
    if content.get('market'):
        add_content_slide(
            f"🎯 {content.get('market_title', 'BOZOR')}",
            content.get('market', ''),
            COLORS['secondary']
        )

    # ==================== 5. BIZNES MODEL ====================
    if content.get('business_model'):
        add_content_slide(
            f"💼 {content.get('business_title', 'BIZNES MODEL')}",
            content.get('business_model', ''),
            COLORS['warning']
        )

    # ==================== 6. RAQOBAT ====================
    if content.get('competition'):
        add_content_slide(
            f"🏆 {content.get('competition_title', 'RAQOBAT')}",
            content.get('competition', ''),
            COLORS['purple']
        )

    # ==================== 7. USTUNLIKLAR ====================
    if content.get('advantage'):
        add_content_slide(
            f"⭐ {content.get('advantage_title', 'USTUNLIKLAR')}",
            content.get('advantage', ''),
            COLORS['primary']
        )

    # ==================== 8. MOLIYA ====================
    if content.get('financials'):
        add_content_slide(
            f"📈 {content.get('financials_title', 'MOLIYA')}",
            content.get('financials', ''),
            COLORS['secondary']
        )

    # ==================== 9. JAMOA ====================
    if content.get('team'):
        add_content_slide(
            f"👥 {content.get('team_title', 'JAMOA')}",
            content.get('team', ''),
            COLORS['accent']
        )

    # ==================== 10. YO'L XARITASI ====================
    if content.get('milestones'):
        add_content_slide(
            f"🗓️ {content.get('milestones_title', 'YOL XARITASI')}",
            content.get('milestones', ''),
            COLORS['warning']
        )

    # ==================== 11. YAKUN ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['accent'], COLORS['secondary'])

    # CTA
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    tf = cta_box.text_frame

    p = tf.add_paragraph()
    p.text = "🚀 KELAJAKNI BIRGALIKDA QURAMIZ!"
    p.font.name = "Calibri"
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    tf.add_paragraph()

    p = tf.add_paragraph()
    p.text = content.get('cta', "Investitsiya qiling va o'sishni kuzating!")
    p.font.name = "Calibri Light"
    p.font.size = Pt(22)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER

    tf.add_paragraph()
    tf.add_paragraph()

    p = tf.add_paragraph()
    p.text = f"📧 {content.get('author', '')}"
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"💼 {content.get('project_name', '')}"
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Faylni saqlash
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"pitch_{package}_{user_id}_{timestamp}.pptx"
    prs.save(filename)

    logger.info(f"Created presentation: {filename}")
    return filename


# ==================== HANDLER'LAR ====================

# VIDEO HANDLER
@dp.callback_query_handler(lambda c: c.data == "watch_video", state='*')
async def watch_video_handler(call: types.CallbackQuery):
    """Video ko'rsatish"""
    await call.answer()

    try:
        await bot.send_video(
            chat_id=call.from_user.id,
            video=VIDEO_FILE_ID,
            caption="📹 Professional pitch tayyorlash bo'yicha qo'llanma\n\n✅ Diqqat bilan ko'ring va amal qiling!"
        )
    except Exception as e:
        logger.error(f"Video send error: {e}")
        await call.message.answer("❌ Video yuborilmadi. Admin bilan bog'laning.")


# START HANDLERS
@dp.callback_query_handler(lambda c: c.data == "start_yes", state='*')
async def start_yes_handler(call: types.CallbackQuery, state: FSMContext):
    """Boshlash tugmasi bosilganda"""
    logger.info(f"start_yes pressed by user {call.from_user.id}")
    await call.answer("Boshlaymiz! ✅")

    user_id = call.from_user.id
    await state.update_data(current_question=0, answers=[])

    text = (
        "📋 Ajoyib! Endi sizga 10 ta savol beraman.\n"
        "Har biriga BATAFSIL javob bering.\n\n"
        "❗ Muhim: Qancha ko'p ma'lumot bersangiz, shuncha yaxshi natija olasiz!\n\n"
        f"{QUESTIONS[0]}"
    )

    await call.message.edit_text(text, reply_markup=cancel_keyboard())
    await PitchStates.waiting_for_answer.set()


@dp.callback_query_handler(lambda c: c.data == "start_no", state='*')
async def start_no_handler(call: types.CallbackQuery):
    """Keyinroq tugmasi bosilganda"""
    logger.info(f"start_no pressed by user {call.from_user.id}")
    await call.answer("Xo'sh, tayyor bo'lganingizda /start ni bosing")
    text = "👌 Mayli, tayyor bo'lganingizda /start buyrug'ini yuboring."
    await call.message.edit_text(text)


# MESSAGE HANDLER
@dp.message_handler(state=PitchStates.waiting_for_answer)
async def answer_handler(message: types.Message, state: FSMContext):
    """Javoblarni ketma-ket qabul qilish"""
    user_id = message.from_user.id
    user_data = await state.get_data()

    current_q = user_data.get('current_question', 0)
    answers = user_data.get('answers', [])

    answers.append(message.text.strip())
    logger.info(f"User {user_id} answered question {current_q + 1}/{len(QUESTIONS)}")

    next_q = current_q + 1

    if next_q < len(QUESTIONS):
        await state.update_data(current_question=next_q, answers=answers)
        progress = f"✅ {next_q}/{len(QUESTIONS)} savol javoblandi\n\n"
        text = progress + QUESTIONS[next_q]
        await message.answer(text, reply_markup=cancel_keyboard())
    else:
        await state.update_data(answers=answers)
        summary = (
            f"🎉 Barcha savollar tugadi!\n\n"
            f"📊 Jami {len(answers)} ta javob qabul qilindi.\n\n"
            f"Endi paketni tanlang:"
        )
        await message.answer(summary, reply_markup=package_keyboard())
        await PitchStates.choosing_package.set()


# Qolgan handler'lar (package, receipt, admin, cancel, status, help, admin panel)
# [Yuqoridagi koddan olinadi, o'zgarmaydi]

# PACKAGE HANDLER
@dp.callback_query_handler(lambda c: c.data.startswith("package_"), state=PitchStates.choosing_package)
async def package_select_handler(call: types.CallbackQuery, state: FSMContext):
    """Paket tanlash"""
    logger.info(f"Package selected by user {call.from_user.id}: {call.data}")
    await call.answer()

    user_id = call.from_user.id
    user_data = await state.get_data()
    answers = user_data.get('answers', [])

    package = "simple" if call.data == "package_simple" else "pro"
    amount = 10000 if package == "simple" else 20000

    order_id = db.create_order(user_id, answers, package)
    logger.info(f"Order {order_id} created for user {user_id}")

    package_name = 'Oddiy' if package == 'simple' else 'Professional'
    formatted_amount = f"{amount:,}"

    payment_text = (
        f"✅ {package_name} paket tanlandi\n\n"
        f"💰 Narx: {formatted_amount} {CURRENCY}\n\n"
        f"💳 TO'LOV MA'LUMOTLARI:\n\n"
        f"Karta raqami:\n"
        f"<code>{CARD_NUMBER}</code>\n\n"
        f"Karta egasi: {CARD_HOLDER}\n"
        f"Summa: <b>{formatted_amount} {CURRENCY}</b>\n\n"
        f"📸 To'lov qilgandan keyin chekni (skrinshot yoki PDF)\n"
        f"bu chatga yuboring.\n\n"
        f"⏳ Admin tasdiqlagach professional PPTX fayli yuboriladi."
    )

    await call.message.edit_text(payment_text, parse_mode="HTML", reply_markup=cancel_keyboard())
    await PitchStates.waiting_for_receipt.set()


# RECEIPT HANDLER
@dp.message_handler(content_types=[ContentType.PHOTO, ContentType.DOCUMENT], state=PitchStates.waiting_for_receipt)
async def receipt_handler(message: types.Message, state: FSMContext):
    """To'lov chekini qabul qilish"""
    user_id = message.from_user.id
    order = db.get_order(user_id)

    if not order:
        await message.reply("❌ Buyurtma topilmadi. Iltimos /start dan boshlang.")
        return

    is_valid, error_msg, file_id = await validate_file(message)
    if not is_valid:
        await message.reply(f"❌ {error_msg}")
        return

    db.update_order(user_id, receipt_file_id=file_id, status="receipt_sent")
    logger.info(f"Receipt uploaded by user {user_id}")

    user = message.from_user
    package_name = 'Oddiy' if order['package'] == 'simple' else 'Professional'
    amount = 10000 if order['package'] == 'simple' else 20000
    formatted_amount = f"{amount:,}"
    username = user.username or 'yoq'
    current_time = datetime.now().strftime('%Y-%m-%d %H:%M')

    info_text = (
        "💳 YANGI TO'LOV CHEKI\n\n"
        "👤 Foydalanuvchi:\n"
        f"  • Ism: {user.full_name}\n"
        f"  • Username: @{username}\n"
        f"  • ID: {user_id}\n\n"
        f"📦 Paket: {package_name}\n"
        f"💰 Summa: {formatted_amount} {CURRENCY}\n\n"
        f"⏰ Vaqt: {current_time}\n"
        f"📝 Javoblar: {len(order['answers'])} ta"
    )

    try:
        if message.content_type == ContentType.PHOTO:
            await bot.send_photo(
                chat_id=ADMIN_ID,
                photo=file_id,
                caption=info_text,
                reply_markup=admin_action_kb(user_id)
            )
        else:
            await bot.send_document(
                chat_id=ADMIN_ID,
                document=file_id,
                caption=info_text,
                reply_markup=admin_action_kb(user_id)
            )

        success_text = (
            "✅ Chek qabul qilindi va adminga yuborildi!\n\n"
            "⏳ Iltimos tasdiqlanishini kuting.\n"
            "Bu odatda 5-30 daqiqa vaqt oladi.\n\n"
            "Holatingizni /status buyrug'i orqali tekshirishingiz mumkin."
        )
        await message.reply(success_text)
        await state.finish()
        logger.info(f"Receipt sent to admin from user {user_id}")

    except Exception as e:
        logger.error(f"Failed to send receipt to admin: {e}")
        error_text = (
            "❌ Adminga yuborishda xatolik yuz berdi.\n"
            "Iltimos keyinroq qayta urinib ko'ring."
        )
        await message.reply(error_text)


# ADMIN ACTIONS
@dp.callback_query_handler(lambda c: c.data.startswith("admin_"), state='*')
async def admin_action_handler(call: types.CallbackQuery):
    """Admin harakatlari"""
    logger.info(f"Admin action by {call.from_user.id}: {call.data}")
    await call.answer()

    if call.from_user.id != ADMIN_ID:
        await call.message.answer("⛔ Sizda ruxsat yo'q!")
        return

    parts = call.data.split(":")
    action = parts[0]
    user_id = int(parts[1])
    order = db.get_order(user_id)

    if not order:
        await call.message.answer("❌ Buyurtma topilmadi")
        return

    if action == "admin_info":
        order_info = format_order_info(order)
        answers_count = len(order['answers'])
        order_id = order['id']

        answers_text = "\n\n📝 JAVOBLAR:\n"
        for i, ans in enumerate(order['answers'], 1):
            answers_text += f"{i}. {ans[:50]}...\n" if len(ans) > 50 else f"{i}. {ans}\n"

        info = (
            "📊 BUYURTMA TAFSILOTLARI\n\n"
            f"{order_info}\n"
            f"📝 Javoblar: {answers_count} ta\n"
            f"🆔 Buyurtma ID: {order_id}"
            f"{answers_text}"
        )
        await call.message.answer(info)
        return

    if action == "admin_reject":
        db.update_order(user_id, status="rejected")
        db.log_admin_action(ADMIN_ID, "reject", user_id)

        try:
            reject_text = (
                "❌ To'lovingiz rad etildi\n\n"
                "Sabablari:\n"
                "• Noto'g'ri summa\n"
                "• Noaniq chek\n"
                "• Boshqa muammo\n\n"
                "Iltimos to'g'ri chekni qayta yuboring yoki /start dan qaytadan boshlang."
            )
            await bot.send_message(chat_id=user_id, text=reject_text)
            await call.message.answer(f"✅ Foydalanuvchi {user_id} ga rad xabari yuborildi")
        except Exception as e:
            logger.error(f"Failed to notify user {user_id}: {e}")
            await call.message.answer(f"❌ Foydalanuvchiga xabar yuborib bo'lmadi: {e}")
        return

    if action == "admin_approve":
        db.update_order(user_id, status="approved")
        db.log_admin_action(ADMIN_ID, "approve", user_id)

        await call.message.answer("⏳ To'lov tasdiqlandi. Professional PPTX tayyorlanmoqda...")

        try:
            answers = order['answers']

            # Mukammal PPTX yaratish
            await call.message.answer("🎨 Professional prezentatsiya yaratilmoqda...")
            pptx_path = await create_stunning_pitch_deck(user_id, answers, order['package'])

            # Foydalanuvchiga yuborish
            package_name = 'Oddiy' if order['package'] == 'simple' else 'Professional'

            caption = (
                "🎉 Sizning professional Pitch Deck tayyor!\n\n"
                f"📦 Paket: {package_name}\n"
                f"✨ AI optimizatsiyasi: {'✅ Maksimal' if USE_OPENAI else '➖ Ochirilgan'}\n"
                f"🌐 Til: O'zbek tili\n"
                f"📄 Slaydlar: 10-12 ta\n"
                f"📊 Bozor tahlillari: ✅ Kiritilgan\n\n"
                "🚀 Investorlarga muvaffaqiyatlar tilaymiz!\n"
                "💡 Maslahat: Prezentatsiyani ko'rib chiqing va kerak bo'lsa tahrirlang."
            )

            with open(pptx_path, "rb") as f:
                await bot.send_document(
                    chat_id=user_id,
                    document=InputFile(f, filename=os.path.basename(pptx_path)),
                    caption=caption
                )

            db.update_order(user_id, status="completed")
            await call.message.answer(f"✅ Professional PPTX foydalanuvchi {user_id} ga yuborildi!")

            # Faylni o'chirish
            if os.path.exists(pptx_path):
                os.remove(pptx_path)
                logger.info(f"Temporary file deleted: {pptx_path}")

        except Exception as e:
            logger.error(f"PPTX generation failed: {e}")
            await call.message.answer(f"❌ Xatolik: {str(e)}")
            db.update_order(user_id, status="error")

            await bot.send_message(
                chat_id=user_id,
                text="❌ Texnik xatolik yuz berdi. Admin bilan bog'laning: @support"
            )


# CANCEL HANDLER
@dp.callback_query_handler(lambda c: c.data == "cancel", state='*')
async def cancel_callback_handler(call: types.CallbackQuery, state: FSMContext):
    """Bekor qilish callback"""
    logger.info(f"Cancel by user {call.from_user.id}")
    await call.answer("Jarayon bekor qilindi")

    current_state = await state.get_state()
    if current_state:
        await state.finish()

    text = (
        "❌ Jarayon bekor qilindi\n\n"
        "Qaytadan boshlash uchun /start buyrug'ini yuboring."
    )
    await call.message.edit_text(text)


@dp.message_handler(commands=['cancel'], state='*')
async def cancel_command_handler(message: types.Message, state: FSMContext):
    """Cancel buyrug'i"""
    current_state = await state.get_state()

    if current_state is None:
        await message.answer("Hozir hech qanday jarayon bajarilmayapti.")
        return

    await state.finish()

    text = (
        "❌ Jarayon bekor qilindi\n\n"
        "Qaytadan boshlash uchun /start buyrug'ini yuboring."
    )
    await message.answer(text)


# COMMAND HANDLERS
@dp.message_handler(commands=['start'], state='*')
async def start_handler(message: types.Message, state: FSMContext):
    """Start command"""
    current_state = await state.get_state()
    if current_state:
        await state.finish()

    user = message.from_user
    db.save_user(user.id, user.username, user.full_name)
    logger.info(f"Start command by user {user.id}")

    text = (
        f"👋 Assalomu alaykum, {user.full_name}!\n\n"
        "🎯 Men sizning startup pitch'ingizni tayyorlashga yordam beraman.\n\n"
        "📝 Jarayon:\n"
        "1️⃣ 10 ta savolga javob bering\n"
        "2️⃣ Paketni tanlang (Oddiy/Pro)\n"
        "3️⃣ Kartaga to'lov qiling\n"
        "4️⃣ Chekni yuboring\n"
        "5️⃣ Professional PPTX oling\n\n"
        "✨ YANGILIKLAR:\n"
        "• O'zbek tilida to'liq content\n"
        "• AI bilan maksimal optimizatsiya\n"
        "• Avtomatik bozor tahlillari\n"
        "• 10+ slayd batafsil ma'lumot bilan\n\n"
        "Boshlaysizmi?"
    )

    await message.answer(text, reply_markup=start_keyboard())


@dp.message_handler(commands=['status'])
async def status_handler(message: types.Message):
    """Buyurtma holatini ko'rish"""
    user_id = message.from_user.id
    order = db.get_order(user_id)

    if not order:
        await message.answer(
            "❌ Buyurtma topilmadi.\n\n"
            "Yangi buyurtma yaratish uchun /start ni bosing."
        )
        return

    order_info = format_order_info(order)
    answers_count = len(order['answers'])

    status = order['status']
    extra_info = ""

    if status == "awaiting_payment":
        extra_info = "\n\n💳 To'lov kutilmoqda. Chekni yuboring."
    elif status == "receipt_sent":
        extra_info = "\n\n⏳ Chek adminga yuborildi. Tasdiq kutilmoqda."
    elif status == "approved":
        extra_info = "\n\n✅ Tasdiqlandi! PPTX tez orada yuboriladi."
    elif status == "rejected":
        extra_info = "\n\n❌ To'lov rad etildi. Qayta yuboring yoki /start."
    elif status == "completed":
        extra_info = "\n\n🎉 Yakunlandi! PPTX yuborilgan."

    text = (
        "📊 BUYURTMA HOLATI\n\n"
        f"{order_info}\n"
        f"📝 Javoblar: {answers_count} ta"
        f"{extra_info}"
    )
    await message.answer(text)


@dp.message_handler(commands=['help'])
async def help_handler(message: types.Message):
    """Yordam"""
    text = (
        "🆘 YORDAM\n\n"
        "📝 Buyruqlar:\n"
        "/start - Yangi pitch boshlash\n"
        "/status - Buyurtma holatini ko'rish\n"
        "/cancel - Jarayonni bekor qilish\n"
        "/help - Yordam\n\n"
        "💡 Qanday ishlaydi?\n"
        "1. /start ni bosing\n"
        "2. 10 ta savolga javob bering\n"
        "3. Paket tanlang\n"
        "4. Kartaga to'lov qiling\n"
        "5. Chekni yuboring\n"
        "6. Admin tasdiqlagach PPTX oling\n\n"
        "❓ Savol: @support"
    )
    await message.answer(text)


# ADMIN COMMANDS
@dp.message_handler(commands=['admin'], user_id=ADMIN_ID)
async def admin_panel_handler(message: types.Message):
    """Admin panel"""
    conn = sqlite3.connect('pitch_bot.db')
    cursor = conn.cursor()

    cursor.execute('SELECT COUNT(*) FROM orders')
    total_orders = cursor.fetchone()[0]

    cursor.execute('SELECT COUNT(*) FROM orders WHERE status = "completed"')
    completed = cursor.fetchone()[0]

    cursor.execute('SELECT COUNT(*) FROM orders WHERE status = "receipt_sent"')
    pending = cursor.fetchone()[0]

    cursor.execute('SELECT COUNT(*) FROM users')
    total_users = cursor.fetchone()[0]

    conn.close()

    text = (
        "👨‍💼 ADMIN PANEL\n\n"
        "📊 Statistika:\n"
        f"  • Jami buyurtmalar: {total_orders}\n"
        f"  • Yakunlangan: {completed}\n"
        f"  • Kutilmoqda: {pending}\n"
        f"  • Foydalanuvchilar: {total_users}\n\n"
        "Yangi chek kelganda sizga xabar beriladi."
    )
    await message.answer(text)


# ERROR HANDLERS
@dp.message_handler(state=PitchStates.waiting_for_receipt)
async def wrong_message_receipt_state(message: types.Message):
    """Receipt kutilayotganda text xabar"""
    await message.reply(
        "⚠️ Iltimos to'lov chekini (rasm yoki PDF) yuboring.\n\n"
        "Agar to'lovni bekor qilmoqchi bo'lsangiz, /cancel ni bosing."
    )


@dp.message_handler(state=PitchStates.choosing_package)
async def wrong_message_package_state(message: types.Message):
    """Paket tanlash state da text xabar"""
    await message.reply(
        "⚠️ Iltimos yuqoridagi tugmalardan paketni tanlang.\n\n"
        "Bekor qilish: /cancel"
    )


@dp.message_handler(content_types=ContentType.ANY, state='*')
async def unknown_content_handler(message: types.Message):
    """Noma'lum content type"""
    await message.reply(
        "❓ Tushunmadim.\n\n"
        "Buyruqlar:\n"
        "/start - Boshlash\n"
        "/status - Holat\n"
        "/help - Yordam"
    )


@dp.errors_handler()
async def errors_handler(update, exception):
    """Global error handler"""
    logger.error(f"Update {update} caused error {exception}")
    return True


# DEBUG
if __name__ == '__main__':
    logger.info("=" * 50)
    logger.info("Bot ishga tushmoqda...")
    logger.info(f"Admin ID: {ADMIN_ID}")
    logger.info(f"OpenAI: {'Enabled' if USE_OPENAI else 'Disabled'}")
    logger.info(f"Database: pitch_bot.db")
    logger.info("=" * 50)