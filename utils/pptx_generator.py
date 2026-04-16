"""
Professional PPTX Generator — Gamma-darajaga yaqin sifat
python-pptx kutubxonasi asosida

Xususiyatlar:
- 11 ta professional rang temalari
- 5 xil slayd layouti (title, standard, card, accent-bar, image+text, xulosa)
- Gradient backgroundlar
- Shadow effektlar
- Smart text autofit
- Pixabay rasm integratsiyasi
- 16:9 format
"""

import logging
import os
import asyncio
import aiohttp
import tempfile
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)

# =====================================================================
#  RANGLAR TEMALARI — har biri to'liq rang sxemasi
# =====================================================================

THEMES = {
    "modern_blue": {
        "name": "Modern Blue",
        "title_bg": ((6, 44, 110), (14, 77, 164)),
        "slide_bg": (248, 250, 252),
        "accent": (14, 107, 247),
        "accent2": (56, 189, 248),
        "title_text": (255, 255, 255),
        "title_on_light": (10, 31, 68),
        "body_text": (51, 64, 85),
        "bullet_accent": (14, 107, 247),
        "subtitle_text": (180, 200, 230),
        "card_bg": (255, 255, 255),
        "card_border": (226, 232, 240),
    },
    "dark_elegant": {
        "name": "Dark Elegant",
        "title_bg": ((18, 18, 25), (40, 40, 60)),
        "slide_bg": (24, 24, 32),
        "accent": (139, 92, 246),
        "accent2": (167, 139, 250),
        "title_text": (255, 255, 255),
        "title_on_light": (240, 240, 255),
        "body_text": (200, 200, 215),
        "bullet_accent": (139, 92, 246),
        "subtitle_text": (150, 150, 180),
        "card_bg": (36, 36, 52),
        "card_border": (55, 55, 75),
    },
    "minimalist": {
        "name": "Minimalist",
        "title_bg": ((245, 245, 245), (235, 235, 235)),
        "slide_bg": (255, 255, 255),
        "accent": (40, 40, 40),
        "accent2": (120, 120, 120),
        "title_text": (25, 25, 25),
        "title_on_light": (25, 25, 25),
        "body_text": (60, 60, 60),
        "bullet_accent": (40, 40, 40),
        "subtitle_text": (100, 100, 100),
        "card_bg": (250, 250, 250),
        "card_border": (220, 220, 220),
    },
    "ocean_fresh": {
        "name": "Ocean Fresh",
        "title_bg": ((0, 80, 100), (0, 128, 128)),
        "slide_bg": (240, 253, 250),
        "accent": (0, 150, 136),
        "accent2": (38, 198, 218),
        "title_text": (255, 255, 255),
        "title_on_light": (0, 77, 64),
        "body_text": (38, 70, 83),
        "bullet_accent": (0, 150, 136),
        "subtitle_text": (178, 223, 219),
        "card_bg": (255, 255, 255),
        "card_border": (204, 239, 233),
    },
    "purple_premium": {
        "name": "Purple Premium",
        "title_bg": ((49, 10, 101), (88, 28, 135)),
        "slide_bg": (250, 245, 255),
        "accent": (147, 51, 234),
        "accent2": (192, 132, 252),
        "title_text": (255, 255, 255),
        "title_on_light": (49, 10, 101),
        "body_text": (59, 47, 85),
        "bullet_accent": (147, 51, 234),
        "subtitle_text": (196, 181, 253),
        "card_bg": (255, 255, 255),
        "card_border": (233, 213, 255),
    },
    "coral_warm": {
        "name": "Coral Warm",
        "title_bg": ((180, 55, 50), (220, 90, 70)),
        "slide_bg": (255, 247, 245),
        "accent": (239, 68, 68),
        "accent2": (251, 146, 60),
        "title_text": (255, 255, 255),
        "title_on_light": (127, 29, 29),
        "body_text": (68, 51, 51),
        "bullet_accent": (239, 68, 68),
        "subtitle_text": (254, 202, 202),
        "card_bg": (255, 255, 255),
        "card_border": (254, 226, 226),
    },
    "rose_creative": {
        "name": "Rose Creative",
        "title_bg": ((157, 23, 77), (190, 24, 93)),
        "slide_bg": (253, 242, 248),
        "accent": (236, 72, 153),
        "accent2": (244, 114, 182),
        "title_text": (255, 255, 255),
        "title_on_light": (131, 24, 67),
        "body_text": (76, 47, 62),
        "bullet_accent": (236, 72, 153),
        "subtitle_text": (251, 207, 232),
        "card_bg": (255, 255, 255),
        "card_border": (252, 231, 243),
    },
    "colorful_bright": {
        "name": "Colorful Bright",
        "title_bg": ((37, 99, 235), (79, 70, 229)),
        "slide_bg": (248, 250, 252),
        "accent": (79, 70, 229),
        "accent2": (16, 185, 129),
        "title_text": (255, 255, 255),
        "title_on_light": (30, 58, 138),
        "body_text": (51, 65, 85),
        "bullet_accent": (79, 70, 229),
        "subtitle_text": (191, 219, 254),
        "card_bg": (255, 255, 255),
        "card_border": (224, 231, 255),
    },
    "warm_classic": {
        "name": "Warm Classic",
        "title_bg": ((120, 80, 40), (160, 110, 60)),
        "slide_bg": (253, 251, 247),
        "accent": (180, 130, 70),
        "accent2": (210, 170, 110),
        "title_text": (255, 255, 255),
        "title_on_light": (80, 50, 20),
        "body_text": (68, 55, 40),
        "bullet_accent": (180, 130, 70),
        "subtitle_text": (220, 200, 170),
        "card_bg": (255, 255, 255),
        "card_border": (237, 228, 210),
    },
    "cosmic_dark": {
        "name": "Cosmic Dark",
        "title_bg": ((10, 5, 30), (30, 15, 60)),
        "slide_bg": (16, 12, 38),
        "accent": (99, 102, 241),
        "accent2": (236, 72, 153),
        "title_text": (255, 255, 255),
        "title_on_light": (230, 230, 255),
        "body_text": (200, 200, 225),
        "bullet_accent": (99, 102, 241),
        "subtitle_text": (150, 150, 200),
        "card_bg": (28, 22, 56),
        "card_border": (50, 42, 80),
    },
    "green_nature": {
        "name": "Green Nature",
        "title_bg": ((6, 78, 59), (21, 128, 61)),
        "slide_bg": (240, 253, 244),
        "accent": (22, 163, 74),
        "accent2": (74, 222, 128),
        "title_text": (255, 255, 255),
        "title_on_light": (6, 78, 59),
        "body_text": (38, 70, 50),
        "bullet_accent": (22, 163, 74),
        "subtitle_text": (187, 247, 208),
        "card_bg": (255, 255, 255),
        "card_border": (209, 250, 229),
    },
}

# Mavjud bot theme ID → yangi tema mapping
THEME_ID_MAP = {
    "chisel": "minimalist",
    "coal": "dark_elegant",
    "blues": "modern_blue",
    "elysia": "rose_creative",
    "breeze": "ocean_fresh",
    "aurora": "purple_premium",
    "coral-glow": "coral_warm",
    "gamma": "colorful_bright",
    "creme": "warm_classic",
    "gamma-dark": "cosmic_dark",
}

DEFAULT_THEME = "modern_blue"

# Slayd o'lchamlari (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =====================================================================
#  GENERATOR CLASS
# =====================================================================

class ProPPTXGenerator:
    """
    Professional PPTX generator — python-pptx asosida

    Foydalanish:
        gen = ProPPTXGenerator(theme_id="blues")
        success = await gen.generate(content, "output.pptx")
    """

    def __init__(self, theme_id: str = None):
        theme_name = DEFAULT_THEME
        if theme_id:
            theme_name = THEME_ID_MAP.get(theme_id.lower(), theme_id.lower())
        if theme_name not in THEMES:
            theme_name = DEFAULT_THEME

        self.theme = THEMES[theme_name]
        self.theme_name = theme_name
        self.prs = None

    # ======================== MAIN API ========================

    async def generate(
        self,
        content: Dict,
        output_path: str,
        pixabay_api_key: str = None,
    ) -> bool:
        """
        Professional PPTX yaratish

        Args:
            content: GPT-4o dan kelgan content dict
            output_path: Chiqish fayl yo'li
            pixabay_api_key: Pixabay API kaliti (ixtiyoriy)

        Returns:
            True — muvaffaqiyatli
        """
        try:
            # 1. Rasmlarni yuklab olish (Unsplash bepul, Pixabay fallback)
            images = await self._fetch_images(content, pixabay_api_key)

            # 2. PPTX yaratish
            self._build(content, images, output_path)

            # 3. Vaqtinchalik rasmlarni tozalash
            for img_path in images.values():
                try:
                    if img_path and os.path.exists(img_path):
                        os.remove(img_path)
                except Exception:
                    pass

            file_size = os.path.getsize(output_path)
            logger.info(f"PPTX yaratildi: {output_path} ({file_size:,} bytes, theme: {self.theme_name})")
            return True

        except Exception as e:
            logger.error(f"PPTX generate xato: {e}", exc_info=True)
            return False

    # ======================== BUILD ========================

    def _build(self, content: Dict, images: Dict, output_path: str):
        """PPTX ni qurib saqlash"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

        title = content.get("title", "Prezentatsiya")
        subtitle = content.get("subtitle", "")
        slides = content.get("slides", [])

        total_slides = len(slides) + 1  # +1 title slayd

        # Title slayd
        self._create_title_slide(title, subtitle)

        # Reja (agenda) slayd — 5+ slayd bo'lgandagina
        if len(slides) >= 5:
            slide_titles = [s.get("title", "") for s in slides]
            self._create_agenda_slide(slide_titles)

        # Content slaydlar
        layout_cycle = [0, 1, 2, 3, 4]  # standard, card, accent-bar, split, stats
        for i, slide_data in enumerate(slides):
            is_last = (i == len(slides) - 1)

            if is_last and len(slides) > 1:
                self._create_conclusion_slide(slide_data)
            else:
                variant = layout_cycle[i % len(layout_cycle)]
                img_path = images.get(i)
                self._create_content_slide(slide_data, variant, img_path)

        # "Rahmat" yakuniy slayd
        self._create_thank_you_slide(title)

        # Slayd raqamlash (title va thank you dan tashqari)
        total_slides = len(self.prs.slides)
        self._add_slide_numbers(total_slides)

        self.prs.save(output_path)

    # ======================== TITLE SLIDE ========================

    def _create_title_slide(self, title: str, subtitle: str):
        """Premium title slayd — gradient bg, geometric dekor, professional branding"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        t = self.theme

        # Gradient background
        self._set_gradient_bg(slide, t["title_bg"][0], t["title_bg"][1])

        # Yuqori chap — dekorativ chiziq
        self._add_rect(slide, Inches(1), Inches(0.6),
                        Inches(3), Inches(0.06), t["accent2"])

        # Yuqori o'ng — kichik dekorativ elementlar
        self._add_rect(slide, Inches(11), Inches(0.5),
                        Inches(0.6), Inches(0.6), t["accent"], alpha=25)
        self._add_rect(slide, Inches(11.8), Inches(0.5),
                        Inches(0.3), Inches(0.3), t["accent2"], alpha=35)

        # Chap vertikal accent chiziq (branding)
        self._add_rect(slide, Inches(0.7), Inches(1.8),
                        Inches(0.08), Inches(3.5), t["accent"])

        # Title
        self._add_textbox(
            slide, title,
            x=Inches(1.2), y=Inches(2.0),
            w=Inches(10.5), h=Inches(2.2),
            font_size=44, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Accent chiziq title ostida
        self._add_rect(slide, Inches(1.2), Inches(4.3),
                        Inches(4.5), Inches(0.08), t["accent2"])

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                x=Inches(1.2), y=Inches(4.7),
                w=Inches(9), h=Inches(1.3),
                font_size=20, bold=False,
                color=t["subtitle_text"],
                alignment=PP_ALIGN.LEFT,
            )

        # Pastki dekorativ zona
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.14),
                        SLIDE_W, Inches(0.14), t["accent"])
        self._add_rect(slide, SLIDE_W - Inches(4), SLIDE_H - Inches(0.14),
                        Inches(4), Inches(0.14), t["accent2"])

        # Pastki chap — kichik geometrik element
        self._add_rect(slide, Inches(1), SLIDE_H - Inches(0.7),
                        Inches(0.35), Inches(0.35), t["accent2"], alpha=30)

    # ======================== CONTENT SLIDES ========================

    def _create_content_slide(self, data: Dict, variant: int, image_path: str = None):
        """Content slayd — layout variant tanlash"""
        has_image = image_path and os.path.exists(image_path)

        if has_image:
            # Rasmli variantlar almashib turadi
            if variant in (1, 3):
                self._create_card_slide(data, image_path)
            else:
                self._create_image_content_slide(data, image_path)
        elif variant == 0:
            self._create_standard_slide(data)
        elif variant == 1:
            self._create_card_slide(data)
        elif variant == 2:
            self._create_accent_bar_slide(data)
        elif variant == 3:
            self._create_split_slide(data)
        elif variant == 4:
            self._create_highlight_slide(data)
        else:
            self._create_standard_slide(data)

    def _create_standard_slide(self, data: Dict):
        """Standard layout — rangli title bar yuqorida, kontent pastda"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title bar (to'liq kenglik)
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.5), t["title_bg"][0])

        # Title bar ostida accent chiziq
        self._add_rect(slide, Inches(0), Inches(1.5),
                        SLIDE_W, Inches(0.06), t["accent"])

        # Title text
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.3),
            w=Inches(11.7), h=Inches(0.9),
            font_size=30, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Content va bullets
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(2.0)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(0.8), y=y,
                w=Inches(11.7), h=Inches(1.6),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.7)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(1.0), y=y,
                w=Inches(11.3), h=SLIDE_H - y - Inches(0.5),
                font_size=16, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

        # Pastki chiziq
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.06),
                        SLIDE_W, Inches(0.06), t["accent"])

    def _create_card_slide(self, data: Dict, image_path: str = None):
        """Card layout — kontent karta ichida, soyali, ixtiyoriy rasm"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title (karta ustida)
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.4),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_on_light"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Title ostida accent chiziq
        self._add_rect(slide, Inches(0.8), Inches(1.3),
                        Inches(3), Inches(0.06), t["accent"])

        # Rasmli yoki rasmsiz karta
        has_image = image_path and os.path.exists(image_path)

        if has_image:
            # Rasm o'ng tomonda, karta chapda
            card = self._add_rounded_rect(
                slide,
                x=Inches(0.6), y=Inches(1.75),
                w=Inches(7.5), h=Inches(5.2),
                fill=t["card_bg"],
                border_color=t.get("card_border"),
                shadow=True,
            )
            self._add_rect(slide, Inches(0.6), Inches(1.75),
                            Inches(0.07), Inches(5.2), t["accent"])

            # Rasm
            img_frame = self._add_rounded_rect(
                slide,
                x=Inches(8.5), y=Inches(1.75),
                w=Inches(4.3), h=Inches(5.2),
                fill=t["card_bg"],
                border_color=t.get("card_border"),
                shadow=True,
            )
            try:
                slide.shapes.add_picture(
                    image_path,
                    left=Inches(8.65), top=Inches(1.9),
                    width=Inches(4.0), height=Inches(4.9),
                )
            except Exception as e:
                logger.warning(f"Rasm qo'shishda xato: {e}")

            content_w = Inches(6.5)
            content_x = Inches(1.3)
            bullet_x = Inches(1.5)
            bullet_w = Inches(6.0)
        else:
            card = self._add_rounded_rect(
                slide,
                x=Inches(0.6), y=Inches(1.75),
                w=Inches(12.1), h=Inches(5.2),
                fill=t["card_bg"],
                border_color=t.get("card_border"),
                shadow=True,
            )
            self._add_rect(slide, Inches(0.6), Inches(1.75),
                            Inches(0.07), Inches(5.2), t["accent"])

            content_w = Inches(10.7)
            content_x = Inches(1.3)
            bullet_x = Inches(1.5)
            bullet_w = Inches(10.2)

        # Card ichida content
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(2.15)

        if content:
            self._add_textbox(
                slide, content,
                x=content_x, y=y,
                w=content_w, h=Inches(1.5),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.7)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=bullet_x, y=y,
                w=bullet_w, h=Inches(3.0),
                font_size=16, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

    def _create_accent_bar_slide(self, data: Dict):
        """Accent bar layout — chapda rangli bar, ikki ustunli bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Chap katta accent bar
        self._add_rect(slide, Inches(0), Inches(0),
                        Inches(0.45), SLIDE_H, t["title_bg"][0])

        # Bar yonida ingichka accent chiziq
        self._add_rect(slide, Inches(0.45), Inches(0),
                        Inches(0.06), SLIDE_H, t["accent"])

        # Yuqori o'ng burchakda kichik dekorativ element
        self._add_rect(slide, SLIDE_W - Inches(2), Inches(0.4),
                        Inches(1.2), Inches(0.06), t["accent2"])

        # Title
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(1.1), y=Inches(0.45),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_on_light"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Divider
        self._add_rect(slide, Inches(1.1), Inches(1.4),
                        Inches(2.5), Inches(0.05), t["accent"])

        # Content
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(1.85)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(1.1), y=y,
                w=Inches(11.7), h=Inches(1.5),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.5)

        # Bullets — 4+ bo'lsa ikki ustunli
        if bullets:
            if len(bullets) >= 4:
                mid = (len(bullets) + 1) // 2
                self._add_bullet_textbox(
                    slide, bullets[:mid],
                    x=Inches(1.1), y=y,
                    w=Inches(5.5), h=SLIDE_H - y - Inches(0.5),
                    font_size=15, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )
                self._add_bullet_textbox(
                    slide, bullets[mid:],
                    x=Inches(7.0), y=y,
                    w=Inches(5.5), h=SLIDE_H - y - Inches(0.5),
                    font_size=15, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )
            else:
                self._add_bullet_textbox(
                    slide, bullets,
                    x=Inches(1.1), y=y,
                    w=Inches(11.7), h=SLIDE_H - y - Inches(0.5),
                    font_size=16, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )

    def _create_image_content_slide(self, data: Dict, image_path: str):
        """Rasm + kontent layout — chapda rasm, o'ngda matn"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title bar
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.4), t["title_bg"][0])
        self._add_rect(slide, Inches(0), Inches(1.4),
                        SLIDE_W, Inches(0.05), t["accent"])

        # Title
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.28),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Rasm (chap tomonda, soyali frame ichida)
        img_frame = self._add_rounded_rect(
            slide,
            x=Inches(0.5), y=Inches(1.8),
            w=Inches(5.8), h=Inches(5.1),
            fill=t["card_bg"],
            border_color=t.get("card_border"),
            shadow=True,
        )

        try:
            slide.shapes.add_picture(
                image_path,
                left=Inches(0.65), top=Inches(1.95),
                width=Inches(5.5), height=Inches(4.8),
            )
        except Exception as e:
            logger.warning(f"Rasm qo'shishda xato: {e}")

        # O'ng tomonda kontent
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(1.85)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(6.7), y=y,
                w=Inches(6.1), h=Inches(2.0),
                font_size=16, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(4.0)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(6.9), y=y,
                w=Inches(5.9), h=SLIDE_H - y - Inches(0.5),
                font_size=15, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

    # ======================== CONCLUSION SLIDE ========================

    def _create_conclusion_slide(self, data: Dict):
        """Xulosa slayd — gradient bg, markazlashtirilgan"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        # Gradient bg (boshqa yo'nalishda)
        self._set_gradient_bg(slide, t["title_bg"][1], t["title_bg"][0], angle=2700000)

        # Yuqori dekorativ chiziq (markazda)
        bar_w = Inches(4)
        bar_x = (SLIDE_W - bar_w) // 2
        self._add_rect(slide, bar_x, Inches(1.8),
                        bar_w, Inches(0.06), t["accent2"])

        # Title
        title = data.get("title", "Xulosa")
        self._add_textbox(
            slide, title,
            x=Inches(1), y=Inches(2.2),
            w=Inches(11.333), h=Inches(1.3),
            font_size=40, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light",
        )

        # Content
        content = data.get("content", "")
        if content:
            self._add_textbox(
                slide, content,
                x=Inches(2), y=Inches(3.8),
                w=Inches(9.333), h=Inches(1.8),
                font_size=19, bold=False,
                color=t["subtitle_text"],
                alignment=PP_ALIGN.CENTER,
                line_spacing=1.5,
            )

        # Bullets (markazda)
        bullets = data.get("bullet_points", [])
        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(3), y=Inches(4.8),
                w=Inches(7.333), h=Inches(2.0),
                font_size=16, color=t["subtitle_text"],
                bullet_color=t["accent2"],
                alignment=PP_ALIGN.CENTER,
            )

        # Pastki accent chiziqlar
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.12),
                        SLIDE_W, Inches(0.12), t["accent"])
        self._add_rect(slide, SLIDE_W - Inches(3.5), SLIDE_H - Inches(0.12),
                        Inches(3.5), Inches(0.12), t["accent2"])

    # ======================== AGENDA SLIDE ========================

    def _create_agenda_slide(self, slide_titles: List[str]):
        """Reja/Mundarija slayd — professional numbered list"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title bar
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.5), t["title_bg"][0])
        self._add_rect(slide, Inches(0), Inches(1.5),
                        SLIDE_W, Inches(0.06), t["accent"])

        # "Reja" sarlavhasi
        self._add_textbox(
            slide, "Reja",
            x=Inches(0.8), y=Inches(0.3),
            w=Inches(11.7), h=Inches(0.9),
            font_size=32, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Slayd sarlavhalarini raqamlangan ro'yxat sifatida ko'rsatish
        # Birinchi va oxirgisini (kirish/xulosa) olib tashlaymiz agar 7+ bo'lsa
        items = slide_titles[:]

        # Ikki ustunli layout (agar 6+ element bo'lsa)
        if len(items) >= 6:
            mid = (len(items) + 1) // 2
            col1 = items[:mid]
            col2 = items[mid:]

            for col_idx, (col_items, col_x) in enumerate([
                (col1, Inches(0.8)),
                (col2, Inches(7.0))
            ]):
                y_pos = Inches(2.0)
                for i, item_title in enumerate(col_items):
                    num = i + 1 if col_idx == 0 else mid + i + 1

                    # Raqam doira
                    circle = self._add_rounded_rect(
                        slide,
                        x=col_x, y=y_pos,
                        w=Inches(0.45), h=Inches(0.45),
                        fill=t["accent"] if (num - 1) % 2 == 0 else t.get("accent2", t["accent"]),
                        shadow=False,
                    )
                    self._add_textbox(
                        slide, str(num),
                        x=col_x, y=y_pos,
                        w=Inches(0.45), h=Inches(0.45),
                        font_size=13, bold=True,
                        color=(255, 255, 255),
                        alignment=PP_ALIGN.CENTER,
                        font_name="Calibri",
                    )

                    # Sarlavha matni
                    self._add_textbox(
                        slide, item_title,
                        x=col_x + Inches(0.65), y=y_pos + Inches(0.02),
                        w=Inches(5.2), h=Inches(0.42),
                        font_size=15, bold=False,
                        color=t["body_text"],
                        alignment=PP_ALIGN.LEFT,
                        font_name="Calibri",
                    )

                    y_pos += Inches(0.65)
        else:
            # Bitta ustunli
            y_pos = Inches(2.0)
            for i, item_title in enumerate(items):
                num = i + 1

                # Raqam doira
                self._add_rounded_rect(
                    slide,
                    x=Inches(1.5), y=y_pos,
                    w=Inches(0.5), h=Inches(0.5),
                    fill=t["accent"] if i % 2 == 0 else t.get("accent2", t["accent"]),
                    shadow=False,
                )
                self._add_textbox(
                    slide, str(num),
                    x=Inches(1.5), y=y_pos,
                    w=Inches(0.5), h=Inches(0.5),
                    font_size=15, bold=True,
                    color=(255, 255, 255),
                    alignment=PP_ALIGN.CENTER,
                    font_name="Calibri",
                )

                # Sarlavha matni
                self._add_textbox(
                    slide, item_title,
                    x=Inches(2.3), y=y_pos + Inches(0.03),
                    w=Inches(9), h=Inches(0.45),
                    font_size=17, bold=False,
                    color=t["body_text"],
                    alignment=PP_ALIGN.LEFT,
                    font_name="Calibri",
                )

                y_pos += Inches(0.72)

        # Pastki chiziq
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.06),
                        SLIDE_W, Inches(0.06), t["accent"])

    # ======================== SPLIT SLIDE ========================

    def _create_split_slide(self, data: Dict):
        """Split layout — chap yarmi gradient bg bilan title, o'ng yarmi content"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Chap panel — gradient
        self._add_rect(slide, Inches(0), Inches(0),
                        Inches(4.8), SLIDE_H, t["title_bg"][0])

        # Accent chiziq — vertikal
        self._add_rect(slide, Inches(4.8), Inches(0),
                        Inches(0.06), SLIDE_H, t["accent"])

        # Title chap panelda
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.6), y=Inches(1.5),
            w=Inches(3.8), h=Inches(2.5),
            font_size=30, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
            line_spacing=1.2,
        )

        # Accent chiziq title ostida
        self._add_rect(slide, Inches(0.6), Inches(4.2),
                        Inches(2), Inches(0.06), t["accent2"])

        # Dekorativ element
        self._add_rect(slide, Inches(0.6), Inches(6.2),
                        Inches(0.4), Inches(0.4), t["accent"], alpha=30)

        # O'ng panel — bullets
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(0.8)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(5.5), y=y,
                w=Inches(7.3), h=Inches(2.2),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.2)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(5.7), y=y,
                w=Inches(7.1), h=SLIDE_H - y - Inches(0.5),
                font_size=15, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

    # ======================== HIGHLIGHT SLIDE ========================

    def _create_highlight_slide(self, data: Dict):
        """Highlight layout — katta raqamli kalit faktlar + pastda bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Yuqorida title bar
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.4), t["title_bg"][0])
        self._add_rect(slide, Inches(0), Inches(1.4),
                        SLIDE_W, Inches(0.06), t["accent"])

        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.28),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Content — markazda kattaroq
        content = data.get("content", "")
        if content:
            # Content ni katta quote shaklida ko'rsatish
            self._add_rect(slide, Inches(0.8), Inches(1.8),
                            Inches(0.08), Inches(1.6), t["accent"])
            self._add_textbox(
                slide, content,
                x=Inches(1.3), y=Inches(1.85),
                w=Inches(11.2), h=Inches(1.6),
                font_size=19, bold=False,
                color=t["body_text"],
                line_spacing=1.5,
            )

        # Bullets — 3 ta ustunli kartochkalar sifatida
        bullets = data.get("bullet_points", [])
        if bullets:
            cols = min(3, len(bullets))
            col_w = Inches(3.5)
            gap = Inches(0.5)
            total_w = col_w * cols + gap * (cols - 1)
            start_x = (SLIDE_W - total_w) // 2

            for idx, bullet in enumerate(bullets[:6]):  # Max 6 ta
                col = idx % cols
                row = idx // cols
                bx = int(start_x) + col * int(col_w + gap)
                by = Inches(3.9) + row * Inches(1.7)

                # Karta
                self._add_rounded_rect(
                    slide,
                    x=bx, y=by,
                    w=col_w, h=Inches(1.45),
                    fill=t["card_bg"],
                    border_color=t.get("card_border"),
                    shadow=True,
                )

                # Raqam
                self._add_textbox(
                    slide, str(idx + 1),
                    x=bx + Inches(0.2), y=by + Inches(0.15),
                    w=Inches(0.5), h=Inches(0.5),
                    font_size=22, bold=True,
                    color=t["accent"],
                    font_name="Calibri Light",
                )

                # Bullet matni
                clean = bullet.strip()
                for prefix in ('•', '●', '○', '▪', '▸', '-', '–', '—', '*'):
                    if clean.startswith(prefix):
                        clean = clean[len(prefix):].strip()
                        break

                self._add_textbox(
                    slide, clean,
                    x=bx + Inches(0.7), y=by + Inches(0.15),
                    w=col_w - Inches(1.0), h=Inches(1.15),
                    font_size=13, bold=False,
                    color=t["body_text"],
                    line_spacing=1.3,
                )

        # Pastki chiziq
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.06),
                        SLIDE_W, Inches(0.06), t["accent"])

    # ======================== THANK YOU SLIDE ========================

    def _create_thank_you_slide(self, title: str):
        """Yakuniy 'Rahmat' slayd — professional branding"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        # Gradient bg
        self._set_gradient_bg(slide, t["title_bg"][0], t["title_bg"][1])

        # Yuqori dekorativ chiziq
        bar_w = Inches(3)
        bar_x = (SLIDE_W - bar_w) // 2
        self._add_rect(slide, bar_x, Inches(2.2),
                        bar_w, Inches(0.06), t["accent2"])

        # "E'tiboringiz uchun rahmat!"
        self._add_textbox(
            slide, "E'tiboringiz uchun rahmat!",
            x=Inches(1), y=Inches(2.6),
            w=Inches(11.333), h=Inches(1.5),
            font_size=42, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light",
        )

        # Prezentatsiya nomi
        self._add_textbox(
            slide, title,
            x=Inches(2), y=Inches(4.3),
            w=Inches(9.333), h=Inches(0.8),
            font_size=18, bold=False,
            color=t["subtitle_text"],
            alignment=PP_ALIGN.CENTER,
        )

        # Pastki accent chiziq
        self._add_rect(slide, bar_x, Inches(5.4),
                        bar_w, Inches(0.06), t["accent2"])

        # Dekorativ elementlar
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.12),
                        SLIDE_W, Inches(0.12), t["accent"])
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.12),
                        Inches(3), Inches(0.12), t["accent2"])

    # ======================== TEXT HELPERS ========================

    def _add_textbox(self, slide, text: str, x, y, w, h,
                     font_size=14, bold=False, color=(0, 0, 0),
                     alignment=PP_ALIGN.LEFT, line_spacing=1.2,
                     font_name="Calibri"):
        """Text box qo'shish — multiline, auto-fit"""
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)

        lines = text.split('\n')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = font_name
            p.font.color.rgb = RGBColor(*color)
            p.alignment = alignment
            p.space_after = Pt(3)
            p.space_before = Pt(1)

            self._set_line_spacing(p, line_spacing)

        # Shrink text on overflow
        self._set_text_autofit(txBox)
        return txBox

    def _add_bullet_textbox(self, slide, bullets: List[str], x, y, w, h,
                            font_size=16, color=(0, 0, 0),
                            bullet_color=(0, 0, 0),
                            alignment=PP_ALIGN.LEFT):
        """Bullet pointlar uchun maxsus text box — har bir bullet alohida paragraf"""
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        for idx, bullet in enumerate(bullets):
            bullet = bullet.strip()
            if not bullet:
                continue

            # Mavjud bullet belgilarini olib tashlash
            for prefix in ('•', '●', '○', '▪', '▸', '-', '–', '—', '*'):
                if bullet.startswith(prefix):
                    bullet = bullet[len(prefix):].strip()
                    break

            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Bullet marker (rangli) + matn
            run_bullet = p.add_run()
            run_bullet.text = "▸  "
            run_bullet.font.size = Pt(font_size)
            run_bullet.font.color.rgb = RGBColor(*bullet_color)
            run_bullet.font.bold = True
            run_bullet.font.name = "Calibri"

            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = RGBColor(*color)
            run_text.font.name = "Calibri"

            p.alignment = alignment
            p.space_after = Pt(8)
            p.space_before = Pt(4)

            self._set_line_spacing(p, 1.3)

        self._set_text_autofit(txBox)
        return txBox

    # ======================== SHAPE HELPERS ========================

    def _add_rect(self, slide, x, y, w, h, fill, alpha=None):
        """To'rtburchak shape qo'shish"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.fill.background()

        if alpha is not None:
            self._set_shape_alpha(shape, alpha)

        return shape

    def _add_rounded_rect(self, slide, x, y, w, h,
                          fill=(255, 255, 255),
                          border_color=None, shadow=False):
        """Rounded rectangle — burchaklari yumaloq"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)

        if border_color:
            shape.line.color.rgb = RGBColor(*border_color)
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()

        if shadow:
            self._add_shadow(shape)

        return shape

    # ======================== BACKGROUND HELPERS ========================

    def _set_gradient_bg(self, slide, color1, color2, angle=5400000):
        """Slide ga gradient background (XML orqali)"""
        cSld = slide._element.find(qn('p:cSld'))
        if cSld is None:
            return

        # Mavjud bg ni olib tashlash
        old_bg = cSld.find(qn('p:bg'))
        if old_bg is not None:
            cSld.remove(old_bg)

        # Yangi bg yaratish va spTree dan OLDIN qo'shish
        bg = etree.Element(qn('p:bg'))
        bgPr = etree.SubElement(bg, qn('p:bgPr'))
        gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        # Color stop 1
        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '0')
        srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
        srgb1.set('val', '%02X%02X%02X' % color1)

        # Color stop 2
        gs2 = etree.SubElement(gsLst, qn('a:gs'))
        gs2.set('pos', '100000')
        srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
        srgb2.set('val', '%02X%02X%02X' % color2)

        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(angle))
        lin.set('scaled', '1')

        etree.SubElement(bgPr, qn('a:effectLst'))

        # bg ni spTree dan oldin qo'shish
        spTree = cSld.find(qn('p:spTree'))
        if spTree is not None:
            cSld.insert(list(cSld).index(spTree), bg)
        else:
            cSld.insert(0, bg)

    def _set_solid_bg(self, slide, color):
        """Slide ga solid background"""
        cSld = slide._element.find(qn('p:cSld'))
        if cSld is None:
            return

        # Mavjud bg ni olib tashlash
        old_bg = cSld.find(qn('p:bg'))
        if old_bg is not None:
            cSld.remove(old_bg)

        # Yangi bg yaratish
        bg = etree.Element(qn('p:bg'))
        bgPr = etree.SubElement(bg, qn('p:bgPr'))
        solidFill = etree.SubElement(bgPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', '%02X%02X%02X' % color)

        etree.SubElement(bgPr, qn('a:effectLst'))

        # bg ni spTree dan oldin qo'shish
        spTree = cSld.find(qn('p:spTree'))
        if spTree is not None:
            cSld.insert(list(cSld).index(spTree), bg)
        else:
            cSld.insert(0, bg)

    # ======================== XML EFFECTS ========================

    def _add_shadow(self, shape):
        """Shape ga drop shadow (XML orqali)"""
        try:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is None:
                return

            effectLst = spPr.find(qn('a:effectLst'))
            if effectLst is None:
                effectLst = etree.SubElement(spPr, qn('a:effectLst'))

            outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
            outerShdw.set('blurRad', '63500')
            outerShdw.set('dist', '25400')
            outerShdw.set('dir', '5400000')
            outerShdw.set('algn', 'tl')
            outerShdw.set('rotWithShape', '0')

            srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
            srgbClr.set('val', '000000')
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '22000')
        except Exception as e:
            logger.debug(f"Shadow xato: {e}")

    def _set_shape_alpha(self, shape, alpha_pct):
        """Shape ga shaffoflik (0-100, 0=to'liq shaffof)"""
        try:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is None:
                return

            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is None:
                return

            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is None:
                return

            a = etree.SubElement(srgbClr, qn('a:alpha'))
            a.set('val', str(alpha_pct * 1000))
        except Exception as e:
            logger.debug(f"Alpha xato: {e}")

    def _set_text_autofit(self, shape):
        """Matn sig'masa avtomatik kichiklaydigan qilish"""
        try:
            txBody = shape._element.find(qn('p:txBody'))
            if txBody is None:
                return

            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is None:
                return

            bodyPr.set('wrap', 'square')
            bodyPr.set('anchor', 't')

            for tag in ['a:noAutofit', 'a:normAutofit', 'a:spAutoFit']:
                existing = bodyPr.find(qn(tag))
                if existing is not None:
                    bodyPr.remove(existing)

            normAutofit = etree.SubElement(bodyPr, qn('a:normAutofit'))
            normAutofit.set('fontScale', '100000')
            normAutofit.set('lnSpcReduction', '20000')
        except Exception as e:
            logger.debug(f"Autofit xato: {e}")

    def _set_line_spacing(self, paragraph, multiplier: float):
        """Paragrafga qator oralig'i qo'yish"""
        try:
            pPr = paragraph._element.find(qn('a:pPr'))
            if pPr is None:
                pPr = etree.SubElement(paragraph._element, qn('a:pPr'))

            # Mavjud lnSpc ni olib tashlash
            existing = pPr.find(qn('a:lnSpc'))
            if existing is not None:
                pPr.remove(existing)

            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(multiplier * 100000)))
        except Exception as e:
            logger.debug(f"Line spacing xato: {e}")

    # ======================== SLIDE NUMBERS ========================

    def _add_slide_numbers(self, total_slides: int):
        """Barcha slaydlarga raqam qo'shish (title va thank you slayddan tashqari)"""
        t = self.theme
        is_dark = self._is_dark_theme()
        last_idx = len(self.prs.slides) - 1

        for idx, slide in enumerate(self.prs.slides):
            if idx == 0 or idx == last_idx:
                continue  # Title va Thank You slaydlarga raqam qo'yilmaydi

            num_text = f"{idx}/{total_slides - 2}"
            # Conclusion slayd uchun subtitle_text, boshqalar uchun body_text
            is_conclusion = idx == last_idx - 1
            num_color = t["subtitle_text"] if is_conclusion else t.get("body_text", (100, 100, 100))

            # Pastki o'ng burchakda kichik raqam
            txBox = slide.shapes.add_textbox(
                SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
                Inches(0.9), Inches(0.3)
            )
            tf = txBox.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = num_text
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(*num_color)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.RIGHT

    def _is_dark_theme(self):
        """Tema qorong'i ekanligini aniqlash"""
        bg = self.theme.get("slide_bg", (255, 255, 255))
        return sum(bg) < 384  # O'rtacha 128 dan past = qorong'i

    # ======================== IMAGE FETCHING ========================

    async def _fetch_images(self, content: Dict, api_key: str = None) -> Dict[int, str]:
        """Rasmlar yuklab olish — Pixabay (asosiy) + Picsum (fallback)"""
        images = {}
        slides = content.get("slides", [])

        timeout = aiohttp.ClientTimeout(total=45)
        # SSL sozlamasi — ba'zi serverlarda sertifikat muammosi bo'lishi mumkin
        try:
            import ssl
            import certifi
            ssl_ctx = ssl.create_default_context(cafile=certifi.where())
            connector = aiohttp.TCPConnector(ssl=ssl_ctx)
        except Exception:
            connector = aiohttp.TCPConnector(ssl=False)

        async with aiohttp.ClientSession(connector=connector, timeout=timeout) as session:
            tasks = []
            for i, slide in enumerate(slides):
                keywords = slide.get("image_keywords", {})
                if keywords:
                    tasks.append(self._fetch_slide_image(session, api_key, i, keywords))

            results = await asyncio.gather(*tasks, return_exceptions=True)
            for result in results:
                if isinstance(result, tuple):
                    idx, path = result
                    if path:
                        images[idx] = path

        logger.info(f"{len(images)} ta rasm yuklab olindi")
        return images

    async def _fetch_slide_image(self, session, api_key: str,
                                  slide_idx: int, keywords: Dict) -> Tuple[int, Optional[str]]:
        """Bitta slayd uchun rasm yuklab olish — Pixabay (asosiy) + Picsum (fallback)"""
        import urllib.parse

        for key_type in ["primary", "secondary", "fallback"]:
            keyword = keywords.get(key_type, "")
            if not keyword:
                continue

            # 1. Pixabay — asosiy (mavzuga mos, professional)
            if api_key:
                img_path = await self._download_pixabay_image(session, api_key, keyword)
                if img_path:
                    return (slide_idx, img_path)

            # 2. Picsum — fallback (har doim ishlaydi, lekin random)
            img_path = await self._download_picsum_image(session, keyword)
            if img_path:
                return (slide_idx, img_path)

        return (slide_idx, None)

    async def _download_pixabay_image(self, session, api_key: str,
                                       keyword: str) -> Optional[str]:
        """Pixabay API dan professional rasm yuklab olish"""
        try:
            import urllib.parse
            encoded_kw = urllib.parse.quote(keyword)
            search_url = (
                f"https://pixabay.com/api/"
                f"?key={api_key}"
                f"&q={encoded_kw}"
                f"&image_type=photo"
                f"&orientation=horizontal"
                f"&per_page=5"
                f"&min_width=800"
                f"&safesearch=true"
            )

            async with session.get(search_url) as resp:
                if resp.status != 200:
                    return None
                data = await resp.json()

            hits = data.get("hits", [])
            if not hits:
                return None

            # Eng katta o'lchamli rasmni tanlash (largeImageURL > webformatURL)
            hit = hits[0]
            img_url = hit.get("largeImageURL") or hit.get("webformatURL", "")
            if not img_url:
                return None

            async with session.get(img_url) as resp:
                if resp.status != 200:
                    return None
                img_data = await resp.read()

            if len(img_data) < 5000:
                return None

            tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            tmp.write(img_data)
            tmp.close()
            logger.debug(f"Pixabay rasm: {keyword} -> {len(img_data)} bytes")
            return tmp.name

        except Exception as e:
            logger.debug(f"Pixabay xato ({keyword}): {e}")
            return None

    async def _download_picsum_image(self, session, keyword: str) -> Optional[str]:
        """Lorem Picsum dan rasm yuklab olish (fallback — har doim ishlaydi)"""
        try:
            import urllib.parse
            import random
            encoded_kw = urllib.parse.quote(keyword)
            img_id = random.randint(1, 1000)
            img_url = f"https://picsum.photos/seed/{encoded_kw}{img_id}/1280/720"

            async with session.get(img_url, allow_redirects=True) as resp:
                if resp.status != 200:
                    return None
                content_type = resp.headers.get('content-type', '')
                if 'image' not in content_type:
                    return None
                img_data = await resp.read()
                if len(img_data) < 5000:
                    return None

            tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            tmp.write(img_data)
            tmp.close()
            logger.debug(f"Picsum rasm: {keyword} -> {len(img_data)} bytes")
            return tmp.name

        except Exception as e:
            logger.debug(f"Picsum xato ({keyword}): {e}")
            return None


# =====================================================================
#  YORDAMCHI FUNKSIYALAR
# =====================================================================

def get_available_themes() -> List[Dict]:
    """Mavjud temalar ro'yxati"""
    result = []
    for theme_id, theme_name in THEME_ID_MAP.items():
        theme = THEMES.get(theme_name, {})
        result.append({
            "id": theme_id,
            "internal_name": theme_name,
            "display_name": theme.get("name", theme_name),
        })
    return result


def resolve_theme_id(theme_id: str) -> str:
    """Bot theme ID ni generator theme nomiga aylantirish"""
    if not theme_id:
        return DEFAULT_THEME
    return THEME_ID_MAP.get(theme_id.lower(), DEFAULT_THEME)
