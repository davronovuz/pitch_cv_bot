# utils/pptx_post_processor.py
# PPTX fayllarni post-processing qilish — shrift, layout, overflow tuzatish
# Presenton API dan olingan PPTX fayllarini professional darajaga keltirish

import logging
import os
from typing import Optional

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.oxml.ns import qn
    from lxml import etree

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx kutubxonasi topilmadi. Post-processing ishlamaydi.")

# Professional dizayn konfiguratsiyasi
STYLE_CONFIG = {
    'font_name': 'Arial',
    'title': {
        'font_size': Pt(28) if PPTX_AVAILABLE else 28,
        'bold': True,
        'color': (0x1A, 0x1A, 0x2E),
        'max_size': Pt(36) if PPTX_AVAILABLE else 36,
        'min_size': Pt(20) if PPTX_AVAILABLE else 20,
    },
    'subtitle': {
        'font_size': Pt(18) if PPTX_AVAILABLE else 18,
        'bold': False,
        'color': (0x55, 0x55, 0x55),
    },
    'body': {
        'font_size': Pt(14) if PPTX_AVAILABLE else 14,
        'bold': False,
        'color': (0x33, 0x33, 0x33),
        'max_size': Pt(18) if PPTX_AVAILABLE else 18,
        'min_size': Pt(10) if PPTX_AVAILABLE else 10,
    },
    'bullet': {
        'font_size': Pt(13) if PPTX_AVAILABLE else 13,
        'bold': False,
        'color': (0x44, 0x44, 0x44),
    },
}


def post_process_pptx(input_path: str, output_path: str = None) -> bool:
    """
    PPTX faylni post-processing qilish — shrift, layout, overflow tuzatish

    Args:
        input_path: Kiruvchi PPTX fayl yo'li
        output_path: Chiquvchi PPTX fayl yo'li (None bo'lsa input_path ga yozadi)

    Returns:
        True — muvaffaqiyatli, False — xato
    """
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx mavjud emas, post-processing o'tkazildi")
        return False

    if output_path is None:
        output_path = input_path

    try:
        prs = Presentation(input_path)
    except Exception as e:
        logger.error(f"PPTX ochib bo'lmadi: {e}")
        return False

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.3)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 1. Layout tuzatish — chegaradan chiqqan elementlarni qaytarish
            _fix_shape_bounds(shape, slide_w, slide_h, margin)

            if not shape.has_text_frame:
                continue

            tf = shape.text_frame

            # 2. Word wrap yoqish
            tf.word_wrap = True

            # 3. Shrift va o'lchamlarni tuzatish
            is_title = _is_title_shape(shape, slide_idx)
            _fix_fonts(tf, is_title)

            # 4. Text overflow tuzatish — XML orqali
            _fix_text_overflow(shape)

            # 5. Margin sozlash
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)

            # 6. XML darajada shriftlarni tuzatish (inherited font fix)
            _fix_fonts_xml(shape._element)

    try:
        prs.save(output_path)
        file_size = os.path.getsize(output_path)
        logger.info(f"Post-processed PPTX saqlandi: {output_path} ({file_size} bytes)")
        return True
    except Exception as e:
        logger.error(f"PPTX saqlashda xato: {e}")
        return False


def _fix_shape_bounds(shape, slide_w, slide_h, margin):
    """Chegaradan chiqqan elementlarni tuzatish"""
    try:
        # Salbiy pozitsiyalarni tuzatish
        if shape.left < 0:
            shape.left = margin
        if shape.top < 0:
            shape.top = margin

        # O'ng chegaradan chiqqanlarni tuzatish
        if shape.left + shape.width > slide_w:
            overflow = (shape.left + shape.width) - slide_w + margin
            if shape.width > overflow:
                shape.width -= overflow
            else:
                shape.left = margin
                shape.width = slide_w - 2 * margin

        # Pastki chegaradan chiqqanlarni tuzatish
        if shape.top + shape.height > slide_h:
            overflow = (shape.top + shape.height) - slide_h + margin
            if shape.height > overflow:
                shape.height -= overflow
            else:
                shape.top = margin
                shape.height = slide_h - 2 * margin

        # Nolga teng o'lchamlarni tuzatish
        if shape.width <= 0:
            shape.width = Inches(2)
        if shape.height <= 0:
            shape.height = Inches(1)
    except Exception as e:
        logger.debug(f"Shape bounds fix xato: {e}")


def _is_title_shape(shape, slide_idx: int) -> bool:
    """Shape title ekanligini aniqlash"""
    try:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 0:
                return True
            if ph_idx == 1 and slide_idx == 0:
                return True
    except Exception:
        pass

    if shape.has_text_frame:
        text = shape.text_frame.text
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and r.font.size >= Pt(24):
                    return True
        if len(text) < 60 and shape.top < Inches(2):
            return True
    return False


def _fix_fonts(text_frame, is_title: bool):
    """Shriftlarni tuzatish va yaxshilash"""
    font_name = STYLE_CONFIG['font_name']

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            # Shrift nomini o'rnatish
            run.font.name = font_name

            if is_title:
                config = STYLE_CONFIG['title']
                # Title shriftini tuzatish
                if run.font.size is None or run.font.size > config['max_size']:
                    run.font.size = config['font_size']
                elif run.font.size < config['min_size']:
                    run.font.size = config['min_size']
                run.font.bold = config['bold']
                run.font.color.rgb = RGBColor(*config['color'])
            else:
                config = STYLE_CONFIG['body']
                # Body shriftini tuzatish
                if run.font.size is None:
                    run.font.size = config['font_size']
                elif run.font.size > config['max_size']:
                    run.font.size = config['max_size']
                elif run.font.size < config['min_size']:
                    run.font.size = config['min_size']
                run.font.color.rgb = RGBColor(*config['color'])

        # Paragraf orasidagi bo'shliq
        paragraph.space_after = Pt(4)
        paragraph.space_before = Pt(2)


def _fix_text_overflow(shape):
    """Matn overflow ni XML orqali tuzatish — shrink text on overflow"""
    try:
        txBody = shape._element.find(qn('p:txBody'))
        if txBody is None:
            return

        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is None:
            return

        # Word wrap yoqish
        bodyPr.set('wrap', 'square')

        # Vertikal alignment — yuqoridan
        bodyPr.set('anchor', 't')

        # Ichki marginlarni o'rnatish (EMU: 1 inch = 914400)
        bodyPr.set('lIns', '91440')   # left: ~0.1"
        bodyPr.set('rIns', '91440')   # right: ~0.1"
        bodyPr.set('tIns', '45720')   # top: ~0.05"
        bodyPr.set('bIns', '45720')   # bottom: ~0.05"

        # Mavjud autofit elementlarini o'chirish
        for tag in ['a:noAutofit', 'a:normAutofit', 'a:spAutoFit']:
            existing = bodyPr.find(qn(tag))
            if existing is not None:
                bodyPr.remove(existing)

        # normAutofit qo'shish — matn sig'masa avtomatik kichiklaydi
        normAutofit = etree.SubElement(bodyPr, qn('a:normAutofit'))
        # fontScale 100% — PowerPoint kerak bo'lganda kichiklaydi
        normAutofit.set('fontScale', '100000')

    except Exception as e:
        logger.debug(f"Text overflow fix xato: {e}")


def _fix_fonts_xml(element):
    """XML darajada shriftlarni tuzatish — inherited fontlar uchun"""
    font_name = STYLE_CONFIG['font_name']

    try:
        txBody = element.find(qn('p:txBody'))
        if txBody is None:
            return

        # Barcha run properties ni tuzatish
        for rPr in txBody.iter(qn('a:rPr')):
            for tag in ['a:latin', 'a:ea', 'a:cs']:
                font_el = rPr.find(qn(tag))
                if font_el is None:
                    font_el = etree.SubElement(rPr, qn(tag))
                font_el.set('typeface', font_name)

        # Default run properties ham tuzatish
        for pPr in txBody.iter(qn('a:pPr')):
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is not None:
                for tag in ['a:latin', 'a:ea', 'a:cs']:
                    font_el = defRPr.find(qn(tag))
                    if font_el is None:
                        font_el = etree.SubElement(defRPr, qn(tag))
                    font_el.set('typeface', font_name)

    except Exception as e:
        logger.debug(f"XML font fix xato: {e}")


def validate_pptx(pptx_path: str) -> dict:
    """PPTX faylni tekshirish — sifat nazorati"""
    if not PPTX_AVAILABLE:
        return {'valid': False, 'errors': ['python-pptx mavjud emas']}

    result = {
        'valid': True,
        'errors': [],
        'warnings': [],
        'info': [],
        'slide_count': 0,
    }

    try:
        prs = Presentation(pptx_path)
        result['slide_count'] = len(prs.slides)

        if len(prs.slides) == 0:
            result['valid'] = False
            result['errors'].append('Prezentatsiyada slayd yo\'q')
            return result

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides):
            has_content = False

            for shape in slide.shapes:
                # Chegaradan chiqqanlarni tekshirish
                if shape.left + shape.width > slide_w * 1.1:
                    result['warnings'].append(
                        f'Slayd {slide_idx + 1}: "{shape.name}" o\'ng chegaradan chiqib ketgan')

                if shape.top + shape.height > slide_h * 1.1:
                    result['warnings'].append(
                        f'Slayd {slide_idx + 1}: "{shape.name}" pastki chegaradan chiqib ketgan')

                if shape.has_text_frame:
                    has_content = True
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.size and run.font.size > Pt(72):
                                result['warnings'].append(
                                    f'Slayd {slide_idx + 1}: Juda katta shrift ({run.font.size})')
                            if run.font.size and run.font.size < Pt(6):
                                result['warnings'].append(
                                    f'Slayd {slide_idx + 1}: Juda kichik shrift ({run.font.size})')

                if shape.shape_type == 13:  # Picture
                    has_content = True

            if not has_content:
                result['warnings'].append(f'Slayd {slide_idx + 1}: Bo\'sh slayd')

        if result['errors']:
            result['valid'] = False

    except Exception as e:
        result['valid'] = False
        result['errors'].append(f'PPTX tekshirishda xato: {e}')

    return result